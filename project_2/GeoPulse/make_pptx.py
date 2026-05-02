"""
GeoPulse — 30-slide PPTX presentation generator
Part 1: Project 1 — Interoperable GIS System  (slides 1-15)
Part 2: Project 2 — GeoPulse SOS              (slides 16-30)

Run:  python make_pptx.py
Out:  GeoPulse_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, copy

# ── Colours — match GUI CSS variables ────────────────────────────────────────
# Light theme palette
BG        = RGBColor(0xf0, 0xf2, 0xf5)   # --bg-app
PANEL     = RGBColor(0xff, 0xff, 0xff)   # --bg-panel
PANEL_ALT = RGBColor(0xf4, 0xf5, 0xf7)   # --bg-pane-head
BORDER    = RGBColor(0xe2, 0xe5, 0xea)   # --border
TEXT      = RGBColor(0x1a, 0x1d, 0x23)   # --text
TEXT_SUB  = RGBColor(0x4a, 0x55, 0x68)   # --text-sub
TEXT_MUTE = RGBColor(0x8a, 0x95, 0xa3)   # --text-muted
ACCENT    = RGBColor(0x4f, 0x6e, 0xf7)   # --accent (indigo)
TEAL      = RGBColor(0x0f, 0xd6, 0xc2)   # --accent-teal
GREEN     = RGBColor(0x10, 0xb9, 0x81)   # --btn-wfs
PURPLE    = RGBColor(0x8b, 0x5c, 0xf6)   # --btn-getmap
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)   # --accent-amber
RED       = RGBColor(0xef, 0x44, 0x44)
WHITE     = RGBColor(0xff, 0xff, 0xff)
BLACK     = RGBColor(0x00, 0x00, 0x00)

# Slide accent bar colors by section
P1_BAR = ACCENT   # project 1 = indigo
P2_BAR = TEAL     # project 2 = teal

# ── Slide size — 16:9 ────────────────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

SCREENSHOTS = os.path.join(
    os.path.dirname(__file__),
    "../../client_final/Output_Screenshots"
)


# ═════════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = line_w
    return shape


def no_line(shape):
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
        italic=False, wrap=True, font="Segoe UI"):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font
    return shape


def mono(slide, text, x, y, w, h, size=9.5, color=TEXT_SUB):
    """Monospace code-style text."""
    return txt(slide, text, x, y, w, h, size=size, color=color,
               font="Courier New", wrap=False)


def heading(slide, title, subtitle=None, bar_color=P1_BAR):
    """Top accent bar + title + optional subtitle."""
    # left colour bar
    bar = no_line(rect(slide, Inches(0), Inches(0),
                        Inches(0.18), Inches(7.5), fill=bar_color))
    # title
    txt(slide, title,
        Inches(0.35), Inches(0.2), Inches(12.78), Inches(0.65),
        size=26, bold=True, color=TEXT)
    # rule under title
    no_line(rect(slide, Inches(0.35), Inches(0.82),
                 Inches(12.48), Inches(0.03), fill=BORDER))
    if subtitle:
        txt(slide, subtitle,
            Inches(0.35), Inches(0.88), Inches(12.78), Inches(0.38),
            size=12, color=TEXT_MUTE)


def card(slide, x, y, w, h, title, bullets,
         title_color=ACCENT, bullet_size=11):
    """White card with title and bullet list."""
    r = rect(slide, x, y, w, h, fill=PANEL, line=BORDER)
    # title strip
    no_line(rect(slide, x, y, w, Inches(0.38), fill=PANEL_ALT, line=BORDER))
    txt(slide, title,
        x + Inches(0.14), y + Inches(0.04), w - Inches(0.28), Inches(0.3),
        size=12, bold=True, color=title_color)
    # bullets
    tf_shape = slide.shapes.add_textbox(
        x + Inches(0.14), y + Inches(0.44),
        w - Inches(0.28), h - Inches(0.54))
    tf = tf_shape.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = ("• " if b and not b.startswith(" ") else "") + b
        run.font.size  = Pt(bullet_size)
        run.font.color.rgb = TEXT_SUB
        run.font.name  = "Segoe UI"


def code_card(slide, code_text, x, y, w, h):
    """Light-gray code block."""
    rect(slide, x, y, w, h,
         fill=PANEL_ALT, line=BORDER)
    tf_shape = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.12),
        w - Inches(0.3), h - Inches(0.24))
    tf = tf_shape.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.strip().split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9.5)
        run.font.name = "Courier New"
        run.font.color.rgb = TEXT


def tag(slide, label, x, y, color=ACCENT, w=Inches(1.3)):
    """Small colour tag/badge."""
    r, g, b = color[0], color[1], color[2]
    bg = RGBColor(min(r + 210, 255), min(g + 210, 255), min(b + 210, 255))
    box = rect(slide, x, y, w, Inches(0.27), fill=bg, line=color, line_w=Pt(1))
    txt(slide, label, x + Inches(0.06), y + Inches(0.02), w - Inches(0.12), Inches(0.23),
        size=9, bold=True, color=color, align=PP_ALIGN.CENTER)


def chip(slide, icon, label, x, y, color=ACCENT):
    """Small icon+text chip on white card."""
    r, g, b = color[0], color[1], color[2]
    bg = RGBColor(min(r + 210, 255), min(g + 210, 255), min(b + 210, 255))
    box = rect(slide, x, y, Inches(2.8), Inches(0.7), fill=bg, line=color, line_w=Pt(1))
    txt(slide, icon + "  " + label,
        x + Inches(0.1), y + Inches(0.14), Inches(2.6), Inches(0.35),
        size=12, bold=True, color=color)


def flow(slide, items, y, bar_color=P1_BAR):
    """Horizontal flow: [(label, sublabel), ...] with arrows between."""
    n   = len(items)
    box_w = Inches(2.4)
    gap   = Inches(0.55)
    total = n * box_w + (n - 1) * gap
    x0    = (W - total) / 2

    for i, (lbl, sub) in enumerate(items):
        bx = x0 + i * (box_w + gap)
        r = rect(slide, bx, y, box_w, Inches(0.8), fill=PANEL, line=bar_color)
        txt(slide, lbl, bx + Inches(0.1), y + Inches(0.06),
            box_w - Inches(0.2), Inches(0.35),
            size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        if sub:
            txt(slide, sub, bx + Inches(0.1), y + Inches(0.44),
                box_w - Inches(0.2), Inches(0.3),
                size=9, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
        if i < n - 1:
            txt(slide, "→",
                bx + box_w + Inches(0.1), y + Inches(0.22),
                Inches(0.35), Inches(0.35),
                size=16, color=bar_color, align=PP_ALIGN.CENTER)


def stat_row(slide, stats, y, bar_color=P1_BAR):
    """Row of (number, label) stat boxes."""
    n  = len(stats)
    bw = Inches(12.33 / n - 0.15)
    for i, (num, lbl) in enumerate(stats):
        bx = Inches(0.35 + i * (12.33 / n))
        r  = rect(slide, bx, y, bw, Inches(1.0), fill=PANEL, line=BORDER)
        txt(slide, num, bx, y + Inches(0.06), bw, Inches(0.55),
            size=28, bold=True, color=bar_color, align=PP_ALIGN.CENTER)
        txt(slide, lbl, bx, y + Inches(0.6), bw, Inches(0.35),
            size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)


def data_table(slide, headers, rows, x, y, w, h, col_widths=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    t = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = cw
    row_h = h // n_rows
    for ri in range(n_rows):
        t.rows[ri].height = row_h
    for ci, hdr in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = hdr
        run.font.size = Pt(11); run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Segoe UI"
    for ri, row_data in enumerate(rows, 1):
        for ci, val in enumerate(row_data):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_ALT if ri % 2 == 0 else PANEL
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.size = Pt(10); run.font.color.rgb = TEXT_SUB
            run.font.name = "Segoe UI"


def screenshot(slide, filename, x, y, w, h):
    """Insert a screenshot if the file exists; grey placeholder otherwise."""
    path = os.path.join(SCREENSHOTS, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        r = rect(slide, x, y, w, h, fill=PANEL_ALT, line=BORDER)
        txt(slide, "[ screenshot: " + filename + " ]",
            x, y + h/2 - Inches(0.2), w, Inches(0.4),
            size=9, color=TEXT_MUTE, align=PP_ALIGN.CENTER)


def divider_slide(title, subtitle, bar_color=P1_BAR):
    """Full-bleed section divider."""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = TEXT   # dark bg for dividers only
    no_line(rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5),
                 fill=bar_color))
    txt(s, title,
        Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.2),
        size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(s, subtitle,
        Inches(0.6), Inches(3.85), Inches(12.1), Inches(0.6),
        size=18, color=RGBColor(0xc0, 0xc8, 0xd8), align=PP_ALIGN.LEFT)
    txt(s, "GNR629  ·  CSRE, IIT Bombay",
        Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
        size=11, color=RGBColor(0x5c, 0x68, 0x82), align=PP_ALIGN.LEFT)
    return s


# ═════════════════════════════════════════════════════════════════════════════
#  PART 1: PROJECT 1 — INTEROPERABLE GIS SYSTEM (slides 1–15)
# ═════════════════════════════════════════════════════════════════════════════

# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = new_slide()
no_line(rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5), fill=P1_BAR))
no_line(rect(s, Inches(0.3), Inches(0), Inches(13.03), Inches(0.08), fill=BORDER))

txt(s, "PROJECT 1",
    Inches(0.6), Inches(1.2), Inches(11.0), Inches(0.5),
    size=13, bold=True, color=P1_BAR)
txt(s, "Interoperable GIS System",
    Inches(0.6), Inches(1.7), Inches(11.0), Inches(1.3),
    size=44, bold=True, color=TEXT)
txt(s, "An OGC-compliant WMS / WFS web client for geospatial data visualisation",
    Inches(0.6), Inches(3.0), Inches(11.0), Inches(0.6),
    size=17, color=TEXT_SUB)
no_line(rect(s, Inches(0.6), Inches(3.75), Inches(4.0), Inches(0.04), fill=BORDER))
txt(s, "GNR629  ·  CSRE, IIT Bombay  ·  Manas Avinashe  ·  2026",
    Inches(0.6), Inches(3.9), Inches(11.0), Inches(0.4),
    size=12, color=TEXT_MUTE)
for i, (lbl, col) in enumerate([("OGC WMS", ACCENT), ("OGC WFS", GREEN),
                                  ("GeoServer", PURPLE), ("OpenLayers", AMBER)]):
    tag(s, lbl, Inches(0.6 + i * 1.5), Inches(4.6), col)


# ── Slide 2: Problem & Motivation ─────────────────────────────────────────────
s = new_slide()
heading(s, "Problem & Motivation", "Why interoperability in GIS?")

card(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(4.0),
     "The Challenge", [
     "Geospatial data is distributed across",
     " many servers and formats",
     "Different vendors use incompatible APIs",
     "No standard way to discover, request,",
     " or render spatial layers remotely",
     "Manual download + re-processing is slow",
], title_color=RED)

card(s, Inches(6.55), Inches(1.35), Inches(6.43), Inches(4.0),
     "The Solution — OGC Web Services", [
     "WMS — Web Map Service",
     " Render layers as images on demand",
     "WFS — Web Feature Service",
     " Retrieve vector features as data",
     "Standard URL operations across all servers",
     "Any compliant client talks to any server",
], title_color=GREEN)

no_line(rect(s, Inches(0.35), Inches(5.55), Inches(12.63), Inches(0.04), fill=BORDER))
txt(s, "Application Domain: Soil Moisture Monitoring — Maharashtra, India",
    Inches(0.35), Inches(5.7), Inches(12.63), Inches(0.4),
    size=12, bold=True, color=TEXT_SUB)
txt(s, "Soil moisture is critical for agricultural productivity, drought assessment, and irrigation planning."
       " This project serves Maharashtra soil moisture layers through GeoServer and visualises them via the OGC client.",
    Inches(0.35), Inches(6.1), Inches(12.63), Inches(0.65),
    size=11, color=TEXT_MUTE)


# ── Slide 3: Technology Stack ──────────────────────────────────────────────────
s = new_slide()
heading(s, "Technology Stack")

for i, (title, items, col) in enumerate([
    ("Server Side", [
        "Apache Tomcat — application server",
        "GeoServer 2.x — OGC WMS/WFS server",
        "PostgreSQL + PostGIS — spatial database",
        "Stores Maharashtra layers (countries,",
        " mh_villages, soil_moisture_georef)",
    ], ACCENT),
    ("Client Side", [
        "HTML5 / CSS3 / Vanilla JavaScript",
        "OpenLayers — interactive mapping library",
        "DOMParser — built-in XML parsing",
        "No build tools, no npm, no framework",
        "Single index.html opens in any browser",
    ], GREEN),
    ("OGC Standards", [
        "WMS 1.1 — GetCapabilities, GetMap,",
        " GetFeatureInfo",
        "WFS 2.0 — GetCapabilities, GetFeature",
        "EPSG:4326 and EPSG:3857 CRS support",
        "XML response parsing + display",
    ], PURPLE),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(1.35), Inches(4.1), Inches(4.9),
         title, items, title_color=col)

stat_row(s, [
    ("4", "JS modules"), ("2", "OGC standards"), ("3", "GeoServer layers"), ("0", "npm dependencies")
], Inches(6.5))


# ── Slide 4: System Architecture ──────────────────────────────────────────────
s = new_slide()
heading(s, "System Architecture")

flow(s, [
    ("Browser", "index.html\nOpenLayers"),
    ("HTTP/S", "GET request\nURL params"),
    ("GeoServer", "Tomcat:8080\nPostGIS backend"),
    ("Response", "image/png\nor XML"),
    ("Render", "OL tile layer\nor vector"),
], Inches(1.9), P1_BAR)

# Three boxes below
for i, (title, items, col) in enumerate([
    ("WMS Flow", [
        "GetCapabilities → parse layer list",
        "User selects layer",
        "GetMap → image URL → OL TileWMS",
        "Image rendered as map overlay",
    ], ACCENT),
    ("WFS Flow", [
        "GetCapabilities → parse feature types",
        "GetFeature → GeoJSON response",
        "OL VectorSource + GeoJSON format",
        "Features rendered on map + table",
    ], GREEN),
    ("Click / Info Flow", [
        "Vector click → attributes from OL",
        "WMS click → GetFeatureInfo request",
        "AbortController cancels in-flight reqs",
        "Attributes shown in parsed data panel",
    ], PURPLE),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(3.15), Inches(4.1), Inches(3.5),
         title, items, title_color=col, bullet_size=10)


# ── Slide 5: OGC Standards — WMS & WFS ────────────────────────────────────────
s = new_slide()
heading(s, "OGC Standards — WMS & WFS",
        "How the client communicates with GeoServer")

card(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(5.4),
     "WMS — Web Map Service", [
     "Returns rendered map images (PNG/JPEG)",
     "",
     "Key operations:",
     " GetCapabilities — list available layers",
     " GetMap — request a map image",
     "   LAYERS, BBOX, SRS, WIDTH, HEIGHT",
     " GetFeatureInfo — query features by pixel",
     "",
     "URL example:",
     " ?SERVICE=WMS&REQUEST=GetMap",
     " &LAYERS=gnr629:countries",
     " &BBOX=68,6,98,38&SRS=EPSG:4326",
     " &WIDTH=800&HEIGHT=600&FORMAT=image/png",
], title_color=ACCENT)

card(s, Inches(6.75), Inches(1.35), Inches(6.23), Inches(5.4),
     "WFS — Web Feature Service", [
     "Returns raw vector feature data",
     "",
     "Key operations:",
     " GetCapabilities — list feature types",
     " GetFeature — download features as",
     "   GeoJSON, GML, or KML",
     "   TYPENAMES, MAXFEATURES, BBOX",
     "",
     "URL example:",
     " ?SERVICE=WFS&REQUEST=GetFeature",
     " &TYPENAMES=gnr629:mh_villages",
     " &OUTPUTFORMAT=application/json",
     " &COUNT=50",
], title_color=GREEN)


# ── Slide 6: GetCapabilities ───────────────────────────────────────────────────
s = new_slide()
heading(s, "GetCapabilities — Service Discovery",
        "The first request — answers: what layers does this server have?")

code_card(s, """\
<!-- WMS GetCapabilities response — parsed by xmlParser.js -->

<WMS_Capabilities version="1.1.1">
  <Service>
    <Title>GeoServer Web Map Service</Title>
  </Service>
  <Capability>
    <Layer>
      <Layer queryable="1">
        <Name>gnr629:countries</Name>
        <Title>World Countries</Title>
        <SRS>EPSG:4326</SRS>
        <LatLonBoundingBox minx="-180" miny="-90" maxx="180" maxy="90"/>
      </Layer>
      <Layer queryable="1">
        <Name>gnr629:mh_villages</Name>
        <Title>Maharashtra Villages</Title>
        <LatLonBoundingBox minx="72.6" miny="15.6" maxx="80.9" maxy="22.0"/>
      </Layer>
    </Layer>
  </Capability>
</WMS_Capabilities>""",
    Inches(0.35), Inches(1.35), Inches(6.8), Inches(5.4))

card(s, Inches(7.35), Inches(1.35), Inches(5.63), Inches(2.6),
     "What the client does with this", [
     "DOMParser parses the raw XML",
     "Extracts layer Name, Title, SRS, BBOX",
     "Populates the layer dropdown",
     "Auto-fills BBOX fields on layer select",
     "Auto-fills SRS field from layer CRS",
], title_color=ACCENT)

card(s, Inches(7.35), Inches(4.15), Inches(5.63), Inches(2.6),
     "Key JS — xmlParser.js", [
     "doc = new DOMParser()",
     "      .parseFromString(xml, 'text/xml')",
     "layers = doc.getElementsByTagName('Layer')",
     "name   = el.querySelector('Name').textContent",
     "bbox   = el.querySelector('[minx]')",
], title_color=PURPLE)


# ── Slide 7: WMS GetMap ────────────────────────────────────────────────────────
s = new_slide()
heading(s, "WMS GetMap — Map Visualisation",
        "Render any server layer as an image overlay on the base map")

screenshot(s, "8. getMap_countries.png",
           Inches(0.35), Inches(1.35), Inches(5.5), Inches(4.3))

card(s, Inches(6.1), Inches(1.35), Inches(6.88), Inches(2.1),
     "GetMap Request Built by ogcRequests.js", [
     "SERVICE=WMS & REQUEST=GetMap",
     "LAYERS=<selected layer>",
     "BBOX=minx,miny,maxx,maxy (EPSG:4326)",
     "SRS=EPSG:4326  WIDTH=800  HEIGHT=600",
     "FORMAT=image/png (or jpeg/tiff)",
], title_color=ACCENT)

card(s, Inches(6.1), Inches(3.65), Inches(6.88), Inches(2.05),
     "OpenLayers Integration", [
     "ol.layer.Image + ol.source.ImageWMS",
     "Layer added to map with auto z-index",
     "Multiple layers stack correctly",
     "Layer checklist panel tracks all layers",
], title_color=GREEN)

screenshot(s, "9. getMap_mh_villages.png",
           Inches(6.1), Inches(5.85), Inches(6.88), Inches(1.4))


# ── Slide 8: WFS GetFeature ────────────────────────────────────────────────────
s = new_slide()
heading(s, "WFS GetFeature — Vector Data",
        "Retrieve actual feature geometry and attributes, not just a rendered image")

screenshot(s, "11. getFeature_countries_258.png",
           Inches(0.35), Inches(1.35), Inches(5.5), Inches(4.3))

card(s, Inches(6.1), Inches(1.35), Inches(6.88), Inches(2.1),
     "GetFeature Request", [
     "SERVICE=WFS & REQUEST=GetFeature",
     "TYPENAMES=<selected layer>",
     "OUTPUTFORMAT=application/json",
     "COUNT=<maxFeatures param>",
     "BBOX=<if set by user>",
], title_color=GREEN)

card(s, Inches(6.1), Inches(3.65), Inches(6.88), Inches(2.05),
     "Result", [
     "GeoJSON features loaded into OL VectorLayer",
     "Features rendered as polygons/points",
     "Attributes displayed in parsed data table",
     "Click a feature → highlight + show info",
], title_color=ACCENT)

screenshot(s, "12. getFeature_countries_50.png",
           Inches(6.1), Inches(5.85), Inches(6.88), Inches(1.4))


# ── Slide 9: GetFeatureInfo / Click Query ─────────────────────────────────────
s = new_slide()
heading(s, "GetFeatureInfo — Click to Query",
        "Click anywhere on the map → retrieve feature attributes at that pixel")

card(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.4),
     "How It Works", [
     "User clicks on the map",
     "",
     "1. If a vector (WFS) feature is at that",
     "   pixel → attributes shown immediately",
     "   from the in-memory GeoJSON",
     "",
     "2. If no vector feature → GetFeatureInfo",
     "   request sent to WMS server:",
     "   REQUEST=GetFeatureInfo",
     "   QUERY_LAYERS=<layer>",
     "   I=<pixel X>  J=<pixel Y>",
     "",
     "AbortController cancels any previous",
     "in-flight request on each new click",
], title_color=ACCENT)

code_card(s, """\
// map.js — WMS click handler

sosMap.on("singleclick", async (evt) => {
  const viewRes = view.getResolution();
  const url = wmsSource.getFeatureInfoUrl(
    evt.coordinate, viewRes, "EPSG:4326",
    { INFO_FORMAT: "application/json",
      FEATURE_COUNT: maxFeatures }
  );

  // cancel previous in-flight request
  if (abortCtrl) abortCtrl.abort();
  abortCtrl = new AbortController();

  const res = await fetch(url, { signal: abortCtrl.signal });
  const json = await res.json();
  // display features in parsed data panel
  displayFeatureInfo(json.features);
});""",
    Inches(6.5), Inches(1.35), Inches(6.48), Inches(5.4))


# ── Slide 10: BBOX Selection ──────────────────────────────────────────────────
s = new_slide()
heading(s, "Bounding Box — Three Ways to Set AOI",
        "Area of interest controls both the map view and the request parameters")

for i, (title, bullets, col, sc) in enumerate([
    ("Method 1: Auto-fill on Layer Select",
     ["Click a layer in the dropdown",
      "BBOX auto-fills from GetCapabilities",
      " LatLonBoundingBox parsed by xmlParser.js",
      "SRS also auto-populated",
     ], ACCENT, "19. bbox_selection.png"),
    ("Method 2: Draw on Map",
     ["Click '✏ Draw on Map' button",
      "Drag rectangle on OpenLayers map",
      "ol.interaction.Draw + createBox()",
      "Coords converted EPSG:3857 → EPSG:4326",
     ], GREEN, "20. bbox_update.png"),
    ("Method 3: Manual + Apply",
     ["Type Min Lon/Lat + Max Lon/Lat directly",
      "Click 'Apply to Map'",
      "Validates ±180° / ±90° ranges",
      "Map view pans and zooms to BBOX",
     ], PURPLE, "21. bbox_apply_to_map.png"),
]):
    bx = Inches(0.35 + i * 4.35)
    card(s, bx, Inches(1.35), Inches(4.1), Inches(2.8), title, bullets, title_color=col)
    screenshot(s, sc, bx, Inches(4.35), Inches(4.1), Inches(2.8))


# ── Slide 11: Layer Management ────────────────────────────────────────────────
s = new_slide()
heading(s, "Layer Management",
        "Stack and control multiple WMS/WFS layers simultaneously")

screenshot(s, "10. soil_moisture_layer_stacked.png",
           Inches(0.35), Inches(1.35), Inches(5.5), Inches(5.4))

card(s, Inches(6.1), Inches(1.35), Inches(6.88), Inches(1.9),
     "Added Layers Panel", [
     "Lists every loaded WMS / WFS layer",
     "WMS / WFS badge shown per layer",
     "Toggle visibility on/off",
     "Adjust opacity with slider",
], title_color=ACCENT)

card(s, Inches(6.1), Inches(3.45), Inches(6.88), Inches(1.9),
     "Z-Index & Stacking", [
     "Most recently added layer on top",
     "WMS layers use auto-incrementing zIndex",
     "WFS vector layers always above WMS images",
     "Clear All Layers button resets map",
], title_color=GREEN)

card(s, Inches(6.1), Inches(5.55), Inches(6.88), Inches(1.2),
     "Tested Layers", [
     "gnr629:countries   ·   gnr629:mh_villages",
     "gnr629:soil_moisture_georef",
], title_color=PURPLE)


# ── Slide 12: XML Parsing ─────────────────────────────────────────────────────
s = new_slide()
heading(s, "XML Parsing — xmlParser.js",
        "DOMParser turns raw XML responses into structured data")

code_card(s, """\
// xmlParser.js — parse WMS GetCapabilities

function parseWMSCapabilities(xmlText) {
  const doc = new DOMParser()
               .parseFromString(xmlText, "text/xml");

  const layerEls = doc.querySelectorAll("Layer > Layer");
  const layers   = [];

  layerEls.forEach(el => {
    const name  = el.querySelector("Name")?.textContent;
    const title = el.querySelector("Title")?.textContent;
    const bbox  = el.querySelector("LatLonBoundingBox");
    const srs   = el.querySelector("SRS")?.textContent;

    if (!name) return;

    layers.push({
      name,  title,  srs,
      minx: parseFloat(bbox?.getAttribute("minx")),
      miny: parseFloat(bbox?.getAttribute("miny")),
      maxx: parseFloat(bbox?.getAttribute("maxx")),
      maxy: parseFloat(bbox?.getAttribute("maxy")),
    });
  });

  return layers;   // → populates layer dropdown
}

// Copy raw XML to response panel
function displayXMLResponse(xmlText) {
  document.getElementById("xmlOutput").textContent = xmlText;
}""",
    Inches(0.35), Inches(1.35), Inches(7.0), Inches(5.6))

card(s, Inches(7.55), Inches(1.35), Inches(5.43), Inches(2.6),
     "DOMParser API", [
     "Built into every browser — no libraries",
     "parseFromString(str, 'text/xml')",
     "Returns a navigable DOM tree",
     "querySelector() works on XML too",
     "textContent gives element text value",
], title_color=ACCENT)

card(s, Inches(7.55), Inches(4.15), Inches(5.43), Inches(2.8),
     "Parsed Data Table", [
     "WFS attributes rendered in a 2-col table",
     "Property | Value rows",
     "WMS parsed layers → layer checklist",
     "GetFeatureInfo JSON → attribute rows",
     "Raw XML always shown in xml panel",
], title_color=GREEN)


# ── Slide 13: UI Features ─────────────────────────────────────────────────────
s = new_slide()
heading(s, "User Interface",
        "Clean, professional design with light/dark themes")

screenshot(s, "4. client_whiteTheme.png",
           Inches(0.35), Inches(1.35), Inches(6.0), Inches(3.5))
screenshot(s, "5. client_darkTheme.png",
           Inches(6.65), Inches(1.35), Inches(6.33), Inches(3.5))

txt(s, "Light Theme", Inches(0.35), Inches(4.9), Inches(6.0), Inches(0.3),
    size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
txt(s, "Dark Theme", Inches(6.65), Inches(4.9), Inches(6.33), Inches(0.3),
    size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)

for i, feat in enumerate([
    "Light / dark theme toggle — saved in localStorage",
    "Resizable map / response panes — drag divider",
    "Collapsible map & response panels",
    "Status dot: idle / loading / active / error",
    "Toast notification system",
]):
    txt(s, "• " + feat,
        Inches(0.35), Inches(5.25 + i * 0.32), Inches(12.63), Inches(0.3),
        size=11, color=TEXT_SUB)


# ── Slide 14: GeoServer Workspace & External Servers ─────────────────────────
s = new_slide()
heading(s, "GeoServer Workspace & External Servers",
        "Local PostGIS layers and remote public WMS services")

screenshot(s, "0. gnr629_workspace.png",
           Inches(0.35), Inches(1.35), Inches(4.2), Inches(3.2))

card(s, Inches(4.75), Inches(1.35), Inches(8.23), Inches(3.2),
     "Local GeoServer (localhost:8080)", [
     "Workspace: gnr629",
     "Layers stored in PostgreSQL/PostGIS:",
     " countries — world boundaries",
     " mh_villages — Maharashtra village polygons",
     " soil_moisture_georef — raster soil data",
     "Served via WMS and WFS endpoints",
], title_color=ACCENT)

screenshot(s, "16. external_wms_server_nasa.png",
           Inches(0.35), Inches(4.75), Inches(4.2), Inches(2.5))

card(s, Inches(4.75), Inches(4.75), Inches(8.23), Inches(2.5),
     "External Public WMS Servers (tested)", [
     "NASA GIBS — gibs.earthdata.nasa.gov",
     "NOAA Weather — opengeo.ncep.noaa.gov",
     "GeoSolutions — gs-stable.geo-solutions.it",
     "Client works with any OGC-compliant server",
], title_color=GREEN)


# ── Slide 15: Project 1 Summary ────────────────────────────────────────────────
s = new_slide()
heading(s, "Project 1 — Summary & Results")

stat_row(s, [
    ("3", "GeoServer layers"),
    ("4", "OGC operations"),
    ("3", "External WMS tested"),
    ("2", "Themes"),
    ("0", "npm dependencies"),
], Inches(1.35))

card(s, Inches(0.35), Inches(2.7), Inches(5.9), Inches(3.9),
     "Requirements Met", [
     "✓ WMS GetCapabilities, GetMap, GetFeatureInfo",
     "✓ WFS GetCapabilities, GetFeature",
     "✓ XML response displayed raw + parsed",
     "✓ Multiple stacked layers with z-index",
     "✓ OSM base layer",
     "✓ External WMS/WFS server support",
     "✓ BBOX coordinate input + draw + apply",
     "✓ Format and size selection",
     "✓ Light/dark theme + resizable panels",
], title_color=GREEN)

card(s, Inches(6.55), Inches(2.7), Inches(6.43), Inches(3.9),
     "Technology Takeaways", [
     "DOMParser replaces any XML library",
     "OpenLayers handles CRS transforms internally",
     "AbortController prevents stale responses",
     "CSS variables make theme switching trivial",
     "StaticFiles — no server needed for client",
     "",
     "Foundation for Project 2:",
     " Same map library, same UI patterns",
     " Same XML parsing approach",
     " Extended with SOS tab + FastAPI",
], title_color=ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
#  PART 2: PROJECT 2 — GeoPulse SOS (slides 16–30)
# ═════════════════════════════════════════════════════════════════════════════

# ── Slide 16: Title — Project 2 ───────────────────────────────────────────────
s = new_slide()
no_line(rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5), fill=P2_BAR))
no_line(rect(s, Inches(0.3), Inches(0), Inches(13.03), Inches(0.08), fill=BORDER))

txt(s, "PROJECT 2",
    Inches(0.6), Inches(1.2), Inches(11.0), Inches(0.5),
    size=13, bold=True, color=P2_BAR)
txt(s, "GeoPulse",
    Inches(0.6), Inches(1.7), Inches(11.0), Inches(1.3),
    size=44, bold=True, color=TEXT)
txt(s, "Extending the GIS portal with an OGC Sensor Observation Service (SOS 2.0)",
    Inches(0.6), Inches(3.0), Inches(11.0), Inches(0.6),
    size=17, color=TEXT_SUB)
no_line(rect(s, Inches(0.6), Inches(3.75), Inches(4.0), Inches(0.04), fill=BORDER))
txt(s, "GNR629  ·  CSRE, IIT Bombay  ·  Manas Avinashe  ·  2026",
    Inches(0.6), Inches(3.9), Inches(11.0), Inches(0.4),
    size=12, color=TEXT_MUTE)
for i, (lbl, col) in enumerate([("OGC SOS 2.0", TEAL), ("SensorML", ACCENT),
                                  ("O&M", GREEN), ("FastAPI", PURPLE), ("SQLite", AMBER)]):
    tag(s, lbl, Inches(0.6 + i * 1.6), Inches(4.6), col, w=Inches(1.45))


# ── Slide 17: What Changed in Project 2 ───────────────────────────────────────
s = new_slide()
heading(s, "From Project 1 → Project 2",
        "GeoPulse adds a complete SOS service on top of the existing OGC client",
        bar_color=P2_BAR)

flow(s, [
    ("Project 1\nclient", "WMS + WFS tab\nOpenLayers"),
    ("+ SOS Tab", "25 sensors\nfilter panel"),
    ("+ FastAPI", "Python backend\nSOS API"),
    ("+ SQLite DB", "17,775 rows\nweather data"),
    ("GeoPulse", "Single portal\nport 8000"),
], Inches(1.7), P2_BAR)

card(s, Inches(0.35), Inches(3.1), Inches(5.9), Inches(3.5),
     "Kept from Project 1", [
     "All WMS / WFS functionality intact",
     "Same OpenLayers map library",
     "Same XML parser, same UI patterns",
     "Same light/dark theme system",
     "Same JS module structure",
], title_color=GREEN)

card(s, Inches(6.55), Inches(3.1), Inches(6.43), Inches(3.5),
     "New in Project 2", [
     "SOS tab with second OpenLayers map",
     "FastAPI Python server (port 8000)",
     "SQLite database with weather observations",
     "SensorML document (DescribeSensor)",
     "Three SOS 2.0 operations",
     "Spatial + temporal + parameter filters",
     "Chart.js visualisations",
], title_color=P2_BAR)


# ── Slide 18: OGC Standards Stack ─────────────────────────────────────────────
s = new_slide()
heading(s, "OGC Standards Stack — Project 2",
        "Three standards work together to describe sensors, observations, and how to query them",
        bar_color=P2_BAR)

for i, (title, items, col) in enumerate([
    ("SOS 2.0 — Sensor Observation Service", [
     "Defines HOW to request sensor data",
     "Three core operations:",
     " GetCapabilities — service discovery",
     " DescribeSensor  — sensor identity",
     " GetObservation  — the actual data",
     "Specifies allowed filter types",
     "Standard URL + XML structure",
    ], TEAL),
    ("O&M — Observations & Measurements", [
     "Defines WHAT a sensor reading is:",
     " featureOfInterest — WHERE",
     " observedProperty  — WHAT",
     " procedure         — WHICH sensor",
     " phenomenonTime    — WHEN",
     " result            — THE VALUE",
     "ISO 19156 international standard",
    ], GREEN),
    ("SensorML 2.0 — Sensor Model Language", [
     "Defines the sensor's IDENTITY:",
     " Unique identifier",
     " Physical location (GML Point)",
     " Output capabilities",
     " Units of measurement",
     "Returned by DescribeSensor operation",
    ], ACCENT),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(1.35), Inches(4.1), Inches(5.6),
         title, items, title_color=col)


# ── Slide 19: System Architecture ─────────────────────────────────────────────
s = new_slide()
heading(s, "System Architecture — GeoPulse",
        "Single FastAPI server handles both the SOS API and frontend delivery",
        bar_color=P2_BAR)

flow(s, [
    ("Browser", "sos.js\nfetch()"),
    ("HTTP GET", "/sos/*\nparams"),
    ("FastAPI", "main.py\nport 8000"),
    ("SQLite", "weather.db\n17,775 rows"),
    ("XML / JSON", "O&M structured\nresponse"),
], Inches(1.8), P2_BAR)

for i, (title, items, col) in enumerate([
    ("FastAPI Endpoints", [
     "GET /  → redirect to /app/index.html",
     "GET /sos/capabilities  → XML",
     "GET /sos/sensor        → SensorML",
     "GET /sos/sensors       → JSON list",
     "GET /sos/observations  → XML / JSON",
     "GET /app/*  → static frontend files",
    ], TEAL),
    ("Data Flow", [
     "sos.js builds URL from form inputs",
     "fetch() sends HTTP GET to FastAPI",
     "FastAPI builds parameterized SQL",
     "SQLite returns matching rows",
     "rows_to_xml() wraps in <Observation>",
     "DOMParser in browser parses response",
    ], ACCENT),
    ("Key Design", [
     "API routes mounted BEFORE StaticFiles",
     "so /sos/* is never shadowed",
     "",
     "SOS_BASE = '' when served by FastAPI",
     "SOS_BASE = 'http://127.0.0.1:8000'",
     "  when opened as file://",
    ], GREEN),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(3.05), Inches(4.1), Inches(3.6),
         title, items, title_color=col, bullet_size=10)


# ── Slide 20: Data Pipeline ────────────────────────────────────────────────────
s = new_slide()
heading(s, "Data Pipeline",
        "Real historical weather data — curated for proper time-series structure",
        bar_color=P2_BAR)

flow(s, [
    ("GlobalWeatherRepository.csv", "138,583 rows\n257 cities · 35 MB"),
    ("process_data.py", "select 25 cities\nkeep full history"),
    ("cleaned_weather.csv", "17,775 rows · 2 MB"),
    ("db_setup.py", "pandas → SQLite"),
    ("weather.db", "1.7 MB · ready"),
], Inches(1.7), P2_BAR)

no_line(rect(s, Inches(0.35), Inches(2.85), Inches(12.63), Inches(1.15),
             fill=RGBColor(0xff, 0xf7, 0xed), line=AMBER))
txt(s, "⚠  Original data problem: df.sample(5000) randomly sampled 5,000 rows across all "
       "257 cities — destroying time-series continuity. "
       "Fix: rewrote process_data.py to explicitly select 25 cities and keep ALL their readings (~710 each).",
    Inches(0.55), Inches(2.94), Inches(12.23), Inches(0.98),
    size=11, color=RGBColor(0x78, 0x35, 0x00))

stat_row(s, [
    ("25", "sensor stations"),
    ("~710", "readings each"),
    ("17,775", "total rows"),
    ("5", "parameters"),
    ("2 years", "time span"),
], Inches(4.3), P2_BAR)

card(s, Inches(0.35), Inches(5.6), Inches(12.63), Inches(1.55),
     "Observed Parameters", [
     "temperature (°C)   ·   humidity (%)   ·   wind_speed (kph)   ·   pressure (mb)   ·   precipitation (mm)",
     "Timestamps: 2024-05-16  →  2026-04-30   ·   sensor_id format: 'Tokyo_Japan'",
], title_color=TEAL, bullet_size=11)


# ── Slide 21: SensorML ────────────────────────────────────────────────────────
s = new_slide()
heading(s, "SensorML — DescribeSensor Operation",
        "GET /sos/sensor  →  the sensor's identity document (SensorML 2.0)",
        bar_color=P2_BAR)

code_card(s, """\
<SensorML xmlns="http://www.opengis.net/sensorml/2.0"
          xmlns:gml="http://www.opengis.net/gml/3.2">

  <member>
    <System>
      <identifier>Tokyo_Japan</identifier>
      <description>Weather sensor in Tokyo, Japan</description>

      <location>
        <gml:Point srsName="EPSG:4326">
          <gml:coordinates>139.69,35.69</gml:coordinates>
        </gml:Point>
      </location>

      <outputs>
        <OutputList>
          <output name="temperature"   unit="Celsius"/>
          <output name="humidity"      unit="%"/>
          <output name="wind_speed"    unit="kph"/>
          <output name="pressure"      unit="mb"/>
        </OutputList>
      </outputs>
    </System>
  </member>
  <!-- 9 more sensors across all continents -->
</SensorML>""",
    Inches(0.35), Inches(1.35), Inches(7.0), Inches(5.6))

card(s, Inches(7.55), Inches(1.35), Inches(5.43), Inches(2.6),
     "Key XML Elements", [
     "identifier — unique sensor ID",
     "  links to sensor_id in database",
     "gml:Point — location, EPSG:4326",
     "  coordinates: lon, lat order",
     "OutputList — declares capabilities",
     "  name + unit per parameter",
], title_color=TEAL)

card(s, Inches(7.55), Inches(4.15), Inches(5.43), Inches(2.8),
     "Our Implementation", [
     "10 sensors described across 5 continents",
     "Served as a static XML file from disk",
     "main.py reads file → Response(xml)",
     "",
     "In full istSOS deployment:",
     "  Dynamic — generated from DB",
     "  Registered via InsertSensor",
], title_color=ACCENT)


# ── Slide 22: GetCapabilities ──────────────────────────────────────────────────
s = new_slide()
heading(s, "SOS GetCapabilities — Service Discovery",
        "GET /sos/capabilities  →  declares all operations, parameters, and offerings",
        bar_color=P2_BAR)

code_card(s, """\
<Capabilities version="2.0.0"
              xmlns:sos="http://www.opengis.net/sos/2.0">
  <ServiceIdentification>
    <Title>GeoPulse Sensor Observation Service</Title>
    <ServiceType>OGC:SOS</ServiceType>
    <ServiceTypeVersion>2.0.0</ServiceTypeVersion>
  </ServiceIdentification>

  <OperationsMetadata>
    <Operation name="GetObservation">
      <DCP><HTTP><Get href="/sos/observations"/></HTTP></DCP>
      <Parameter name="param">
        <AllowedValues>
          <Value>temperature</Value>
          <Value>humidity</Value>
          <Value>wind_speed</Value>
          <Value>pressure</Value>
          <Value>precipitation</Value>
        </AllowedValues>
      </Parameter>
    </Operation>
  </OperationsMetadata>

  <ObservationOfferingList>
    <ObservationOffering>
      <identifier>weather_global</identifier>
      <responseFormat>application/xml</responseFormat>
      <responseFormat>application/json</responseFormat>
    </ObservationOffering>
  </ObservationOfferingList>
</Capabilities>""",
    Inches(0.35), Inches(1.35), Inches(7.0), Inches(5.6))

card(s, Inches(7.55), Inches(1.35), Inches(5.43), Inches(2.5),
     "Three Sections", [
     "ServiceIdentification — who we are",
     "OperationsMetadata — what we support",
     "  + allowed filter parameters",
     "ObservationOfferingList — the datasets",
     "  response formats available",
], title_color=TEAL) if False else \
card(s, Inches(7.55), Inches(1.35), Inches(5.43), Inches(2.5),
     "Three Sections", [
     "ServiceIdentification — who we are",
     "OperationsMetadata — what we support",
     "  + allowed filter parameters",
     "ObservationOfferingList — the datasets",
     "  response formats available",
], title_color=TEAL)

card(s, Inches(7.55), Inches(4.05), Inches(5.43), Inches(2.9),
     "Why It Matters", [
     "A client can auto-discover this service",
     "without any prior knowledge of the API",
     "",
     "Standard structure means any OGC tool",
     "(QGIS, ArcGIS, custom client) can",
     "parse it and start querying",
], title_color=ACCENT)


# ── Slide 23: GetObservation Filters ──────────────────────────────────────────
s = new_slide()
heading(s, "GetObservation — Three Filter Dimensions",
        "GET /sos/observations  —  all filters optional and combinable",
        bar_color=P2_BAR)

for i, (title, rows, col) in enumerate([
    ("Spatial", [
        ("country", "Japan", "case-insensitive exact match"),
        ("bbox",    "60,10,140,50", "minLon,minLat,maxLon,maxLat"),
    ], TEAL),
    ("Temporal", [
        ("after", "2025-01-01 00:00:00", "everything after this instant"),
        ("start", "2025-01-01",          "start of time window"),
        ("end",   "2025-03-31",          "end of time window"),
    ], ACCENT),
    ("Parameter", [
        ("param",  "temperature",  "temperature/humidity/wind_speed/..."),
        ("op",     "between",      "eq neq lt lte gt gte between"),
        ("value",  "30",           "comparison value"),
        ("value2", "40",           "upper bound (between only)"),
    ], GREEN),
]):
    bx = Inches(0.35 + i * 4.35)
    card(s, bx, Inches(1.35), Inches(4.1), Inches(3.8),
         title, [f"{k}={v}" + f"\n  {desc}" for k, v, desc in rows],
         title_color=col, bullet_size=10)

no_line(rect(s, Inches(0.35), Inches(5.35), Inches(12.63), Inches(1.3),
             fill=RGBColor(0xf0, 0xf9, 0xf8), line=TEAL))
txt(s, "Example — Asian sensors, temperature 30–40°C, January 2025:",
    Inches(0.55), Inches(5.42), Inches(12.23), Inches(0.3),
    size=10, bold=True, color=TEAL)
txt(s, "/sos/observations?bbox=60,5,150,55&param=temperature&op=between&value=30&value2=40&start=2025-01-01&end=2025-01-31",
    Inches(0.55), Inches(5.75), Inches(12.23), Inches(0.7),
    size=10, color=TEXT, font="Courier New")


# ── Slide 24: SQL Engine ───────────────────────────────────────────────────────
s = new_slide()
heading(s, "Dynamic SQL — main.py Filter Engine",
        "URL parameters become safe, parameterized SQLite queries",
        bar_color=P2_BAR)

code_card(s, """\
conditions  = []
bind_params = []

# Spatial — bounding box
if bbox:
    min_lon, min_lat, max_lon, max_lat = map(float, bbox.split(","))
    conditions.append("longitude BETWEEN ? AND ?")
    conditions.append("latitude  BETWEEN ? AND ?")
    bind_params += [min_lon, max_lon, min_lat, max_lat]

# Spatial — country
if country:
    conditions.append("LOWER(country) = LOWER(?)")
    bind_params.append(country)

# Temporal — time window
if start:
    conditions.append("timestamp >= ?")
    bind_params.append(start)
if end:
    conditions.append("timestamp <= ?")
    bind_params.append(end)

# Parameter filter — param validated against allowlist
ALLOWED = {"temperature","humidity","wind_speed","pressure","precipitation"}
if param in ALLOWED and op == "between":
    conditions.append(f"{param} BETWEEN ? AND ?")
    bind_params += [value, value2]

where = "WHERE " + " AND ".join(conditions) if conditions else ""
sql   = f"SELECT * FROM observations {where} LIMIT ?"
rows  = conn.execute(sql, bind_params).fetchall()""",
    Inches(0.35), Inches(1.35), Inches(8.0), Inches(5.6))

card(s, Inches(8.55), Inches(1.35), Inches(4.43), Inches(2.6),
     "SQL Injection Safety", [
     "? placeholders keep values",
     "  separate from the SQL string",
     "User input never concatenated",
     "  directly into SQL",
     "param column name validated",
     "  against an explicit allowlist",
], title_color=RED)

card(s, Inches(8.55), Inches(4.15), Inches(4.43), Inches(2.8),
     "O&M XML Output", [
     "rows_to_xml() wraps each row:",
     "<Observations count='N'>",
     "  <Observation>",
     "    <temperature>24.0</temperature>",
     "    <timestamp>2025-01-15</timestamp>",
     "    <sensor_id>Tokyo_Japan</sensor_id>",
     "  </Observation>",
], title_color=TEAL)


# ── Slide 25: O&M Mapping ─────────────────────────────────────────────────────
s = new_slide()
heading(s, "O&M Concept Mapping",
        "Every SQLite row is an informal O&M Observation",
        bar_color=P2_BAR)

data_table(s,
    ["O&M Concept", "Definition", "Our SQLite Column"],
    [
        ["featureOfInterest", "The place being observed",
         "location_name, latitude, longitude"],
        ["observedProperty", "What is being measured",
         "column name: temperature / humidity / wind_speed / ..."],
        ["procedure", "The sensor that made the observation",
         "sensor_id  →  links to SensorML <System> entry"],
        ["phenomenonTime", "When the observation was made",
         "timestamp  (e.g. '2025-01-15 11:30:00')"],
        ["result", "The measured value + unit",
         "numeric value + unit declared in SensorML outputs"],
    ],
    Inches(0.35), Inches(1.35), Inches(12.63), Inches(3.6),
    col_widths=[Inches(2.5), Inches(3.8), Inches(5.8)],
)

no_line(rect(s, Inches(0.35), Inches(5.15), Inches(12.63), Inches(1.95),
             fill=RGBColor(0xf0, 0xf9, 0xf8), line=TEAL))
txt(s, "XML response structure from /sos/observations:",
    Inches(0.55), Inches(5.22), Inches(6.0), Inches(0.3),
    size=10, bold=True, color=TEAL)
txt(s, '<Observations count="30">  <Observation>  '
       '<location_name>Berlin</location_name>  <temperature>6.3</temperature>  '
       '<humidity>75</humidity>  <timestamp>2025-01-01 11:30:00</timestamp>  '
       '<sensor_id>Berlin_Germany</sensor_id>  </Observation>  ...',
    Inches(0.55), Inches(5.56), Inches(12.23), Inches(1.4),
    size=9.5, color=TEXT, font="Courier New")


# ── Slide 26: Frontend — SOS Tab ──────────────────────────────────────────────
s = new_slide()
heading(s, "Frontend — SOS Tab Structure",
        "sos.js handles map, fetch pipeline, XML parsing, table, and charts",
        bar_color=P2_BAR)

data_table(s,
    ["File", "Role"],
    [
        ["sos.js (624 lines)", "Entire SOS tab: map init, fetch, XML parse, table, charts, bidirectional sync"],
        ["tabs.js",           "Switches OGC ↔ SOS tab, initialises sosMap on first open (needs div visible)"],
        ["index.html #sosTab","Left filter panel + right sos-map-pane + sos-bottom (table + charts)"],
        ["style.css",         "SOS-specific layout: .sos-body, .sos-panel, .sos-map-pane, .sos-bottom"],
    ],
    Inches(0.35), Inches(1.35), Inches(12.63), Inches(2.3),
    col_widths=[Inches(2.5), Inches(9.6)],
)

card(s, Inches(0.35), Inches(3.85), Inches(5.9), Inches(3.3),
     "sos.js Module State", [
     "sosMap         — OL map instance",
     "sosMarkerSource — vector features (dots)",
     "_currentObs[]  — last fetched observations",
     "_sensorMeta{}  — sensor_id → lat/lon/name",
     "_selectedSensor — currently clicked sensor",
     "_barChart / _pieChart — Chart.js instances",
], title_color=TEAL)

card(s, Inches(6.55), Inches(3.85), Inches(6.43), Inches(3.3),
     "Fetch → Render Pipeline", [
     "1. fetchObservations() reads all form fields",
     "2. Builds URL with URLSearchParams",
     "3. fetch() → raw XML text from FastAPI",
     "4. _parseObsXML() → JS array of objects",
     "5. Stored in _currentObs[]",
     "6. _renderObsMarkers() recolours dots",
     "7. _renderTable() builds <tr> rows",
     "8. _renderCharts() draws bar + doughnut",
], title_color=ACCENT)


# ── Slide 27: SOS Map ──────────────────────────────────────────────────────────
s = new_slide()
heading(s, "SOS Tab — Interactive Map",
        "25 sensor markers coloured by average temperature of loaded observations",
        bar_color=P2_BAR)

for i, (title, items, col) in enumerate([
    ("Temperature Colour Scale", [
     "< 0 °C      Blue  — freezing",
     "0 – 15 °C   Teal  — cool",
     "15 – 25 °C  Green — mild",
     "25 – 35 °C  Orange — warm",
     "> 35 °C     Red   — hot",
     "",
     "Grey = no data loaded yet",
     "Marker size: 7px / 11px (selected)",
    ], TEAL),
    ("Map Layer Stack", [
     "Layer 0: OSM base tiles",
     " ol.layer.Tile + ol.source.OSM",
     "Layer 1: sensor marker dots",
     " ol.layer.Vector  zIndex=10",
     "Layer 2: bbox draw rectangle",
     " ol.layer.Vector  zIndex=20",
     "",
     "All layers in one ol.Map on #sosMap",
    ], ACCENT),
    ("BBOX Draw Tool", [
     "Click '✏ Draw' button",
     "ol.interaction.Draw added to map",
     "Drag rectangle on map",
     "drawend event fires:",
     " geometry.getExtent()",
     " transformExtent → EPSG:4326",
     " fills four bbox input fields",
    ], GREEN),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(1.35), Inches(4.1), Inches(4.6),
         title, items, title_color=col)

card(s, Inches(0.35), Inches(6.15), Inches(12.63), Inches(1.0),
     "Popup Overlay (OL Overlay)", [
     "Attached to document.body so OL absolute positioning works correctly  ·  "
     "Shows: city, country, obs count, latest temp/humidity/wind/pressure/precip/timestamp",
], title_color=TEXT_MUTE, bullet_size=10)


# ── Slide 28: Table & Charts ───────────────────────────────────────────────────
s = new_slide()
heading(s, "SOS Tab — Data Table & Visualisations",
        "Observations rendered as a sortable table and two Chart.js charts",
        bar_color=P2_BAR)

card(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(2.5),
     "Observation Table", [
     "Columns: #, City, Country, Timestamp,",
     "  Temp °C, Humidity %, Wind kph,",
     "  Pressure mb, Precip mm",
     "Up to 500 rows displayed",
     "Temp cells: coloured badge matching legend",
], title_color=ACCENT)

card(s, Inches(6.75), Inches(1.35), Inches(6.23), Inches(2.5),
     "XML Response Panel", [
     "Raw XML from server shown in full",
     "Displays actual O&M-structured response",
     "Request URL shown above — copy-pasteable",
     "⧉ Copy button for clipboard export",
], title_color=TEXT_MUTE)

card(s, Inches(0.35), Inches(4.1), Inches(6.1), Inches(3.05),
     "Bar Chart — Avg Temperature by City", [
     "Chart.js type: 'bar'",
     "Top 10 cities by average temperature",
     "Bar colour = same temperature scale",
     "  as map markers (_tempColor function)",
     "Responsive, maintainAspectRatio: false",
     "Redrawn on every fetch (old instance",
     "  destroyed first to avoid memory leak)",
], title_color=TEAL)

card(s, Inches(6.75), Inches(4.1), Inches(6.23), Inches(3.05),
     "Doughnut Chart — Observation Count", [
     "Chart.js type: 'doughnut'",
     "Shows which sensors have most data",
     "  in the current filter result",
     "Top 8 cities shown in legend",
     "Legend positioned on right side",
     "borderColor matches bg for segment gap",
], title_color=GREEN)


# ── Slide 29: Bidirectional Sync ───────────────────────────────────────────────
s = new_slide()
heading(s, "Bidirectional Marker ↔ Table Sync",
        "Clicking anything makes everything else react — map, table, and popup stay in sync",
        bar_color=P2_BAR)

card(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(3.3),
     "Click a table row  →", [
     "Row highlighted (blue outline)",
     "Map pans + zooms to that sensor",
     "Marker enlarges (radius 7 → 11px)",
     "White border added to marker",
     "Popup opens with latest values",
], title_color=ACCENT)

card(s, Inches(6.55), Inches(1.35), Inches(6.43), Inches(3.3),
     "Click a map marker  →", [
     "Marker enlarges + white stroke",
     "OL Overlay popup anchored to coordinate",
     "Matching table row highlighted",
     "Table auto-scrolls to that row",
     "Previous selection deselected",
], title_color=TEAL)

code_card(s, """\
// Table row click → map pan
tr.addEventListener("click", () => {
  _selectSensor(sensorId);
  const f = sosMarkerSource.getFeatures()
              .find(f => f.get("sensor_id") === sensorId);
  const coord = f.getGeometry().getCoordinates();
  sosMap.getView().animate({ center: coord, zoom: 4, duration: 500 });
  _showSensorPopup(sensorId, coord);
});

// Map marker click → table highlight
sosMap.on("click", (evt) => {
  sosMap.forEachFeatureAtPixel(evt.pixel, (feature) => {
    const sid = feature.get("sensor_id");
    _selectSensor(sid);
    _highlightTableRow(sid);          // adds .sos-row-selected CSS class
    row.scrollIntoView({ behavior: "smooth", block: "nearest" });
  }, { hitTolerance: 8 });
});""",
    Inches(0.35), Inches(4.85), Inches(12.63), Inches(2.4))


# ── Slide 30: Summary & Standards Compliance ──────────────────────────────────
s = new_slide()
heading(s, "Summary — Standards Compliance & Results",
        "GeoPulse implements OGC WMS, WFS, and SOS across both projects",
        bar_color=P2_BAR)

data_table(s,
    ["Standard", "Operation", "Implementation"],
    [
        ["WMS 1.1",     "GetCapabilities, GetMap, GetFeatureInfo",
         "Project 1 — full WMS client against GeoServer"],
        ["WFS 2.0",     "GetCapabilities, GetFeature",
         "Project 1 — WFS client with BBOX + max features"],
        ["SOS 2.0",     "GetCapabilities",
         "GET /sos/capabilities → Capabilities XML"],
        ["SOS 2.0",     "DescribeSensor",
         "GET /sos/sensor → SensorML 2.0 document"],
        ["SOS 2.0",     "GetObservation",
         "GET /sos/observations → spatial + temporal + parameter filters"],
        ["O&M",         "Observation model",
         "Each DB row encodes all 5 O&M components"],
        ["SensorML 2.0","System description",
         "10 sensors with GML location + output list + units"],
    ],
    Inches(0.35), Inches(1.35), Inches(12.63), Inches(4.0),
    col_widths=[Inches(1.5), Inches(3.2), Inches(7.4)],
)

card(s, Inches(0.35), Inches(5.6), Inches(6.1), Inches(1.55),
     "Run Commands", [
     "cd project_2/GeoPulse/backend",
     "pip install -r requirements.txt",
     "uvicorn main:app --host 127.0.0.1 --port 8000",
     "→ open http://127.0.0.1:8000",
], title_color=TEAL, bullet_size=10)

card(s, Inches(6.75), Inches(5.6), Inches(6.23), Inches(1.55),
     "Possible Extensions", [
     "Migrate SOS backend to istSOS + Docker",
     "Real-time data via InsertObservation",
     "WCS (Web Coverage Service) for rasters",
], title_color=ACCENT, bullet_size=10)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "GeoPulse_Presentation.pptx")
prs.save(out)
print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")
