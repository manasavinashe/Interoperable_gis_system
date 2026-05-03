"""
GeoPulse v2 — 30-slide PPTX presentation generator
Part 1: Project 1 — Interoperable GIS System        (slides 1-15)
Part 2: Project 2 — GeoPulse istSOS + Docker + PostGIS (slides 16-30)

Run:  python make_pptx_istsos.py
Out:  GeoPulse_istSOS_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colours — match GUI CSS variables ────────────────────────────────────────
BG        = RGBColor(0xf0, 0xf2, 0xf5)
PANEL     = RGBColor(0xff, 0xff, 0xff)
PANEL_ALT = RGBColor(0xf4, 0xf5, 0xf7)
BORDER    = RGBColor(0xe2, 0xe5, 0xea)
TEXT      = RGBColor(0x1a, 0x1d, 0x23)
TEXT_SUB  = RGBColor(0x4a, 0x55, 0x68)
TEXT_MUTE = RGBColor(0x8a, 0x95, 0xa3)
ACCENT    = RGBColor(0x4f, 0x6e, 0xf7)   # indigo
TEAL      = RGBColor(0x0f, 0xd6, 0xc2)
GREEN     = RGBColor(0x10, 0xb9, 0x81)
PURPLE    = RGBColor(0x8b, 0x5c, 0xf6)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
RED       = RGBColor(0xef, 0x44, 0x44)
WHITE     = RGBColor(0xff, 0xff, 0xff)
BLACK     = RGBColor(0x00, 0x00, 0x00)
DOCKER    = RGBColor(0x08, 0x6d, 0xd7)   # Docker blue
PG_BLUE   = RGBColor(0x33, 0x68, 0xa0)   # PostgreSQL blue

P1_BAR = ACCENT
P2_BAR = TEAL

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

SCREENSHOTS = os.path.join(
    os.path.dirname(__file__),
    "../../client_final/Output_Screenshots"
)


# ═════════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = line_w
    return shape


def no_line(shape):
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
        italic=False, wrap=True, font="Segoe UI"):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return shape


def mono(slide, text, x, y, w, h, size=9.5, color=TEXT_SUB):
    return txt(slide, text, x, y, w, h, size=size, color=color,
               font="Courier New", wrap=False)


def heading(slide, title, subtitle=None, bar_color=P1_BAR):
    no_line(rect(slide, Inches(0), Inches(0),
                 Inches(0.18), Inches(7.5), fill=bar_color))
    txt(slide, title,
        Inches(0.35), Inches(0.2), Inches(12.78), Inches(0.65),
        size=26, bold=True, color=TEXT)
    no_line(rect(slide, Inches(0.35), Inches(0.82),
                 Inches(12.48), Inches(0.03), fill=BORDER))
    if subtitle:
        txt(slide, subtitle,
            Inches(0.35), Inches(0.88), Inches(12.78), Inches(0.38),
            size=12, color=TEXT_MUTE)


def card(slide, x, y, w, h, title, bullets,
         title_color=ACCENT, bullet_size=11):
    rect(slide, x, y, w, h, fill=PANEL, line=BORDER)
    no_line(rect(slide, x, y, w, Inches(0.38), fill=PANEL_ALT, line=BORDER))
    txt(slide, title,
        x + Inches(0.14), y + Inches(0.04), w - Inches(0.28), Inches(0.3),
        size=12, bold=True, color=title_color)
    tf_shape = slide.shapes.add_textbox(
        x + Inches(0.14), y + Inches(0.44),
        w - Inches(0.28), h - Inches(0.54))
    tf = tf_shape.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = ("• " if b and not b.startswith(" ") else "") + b
        run.font.size  = Pt(bullet_size)
        run.font.color.rgb = TEXT_SUB
        run.font.name  = "Segoe UI"


def code_card(slide, code_text, x, y, w, h):
    rect(slide, x, y, w, h, fill=PANEL_ALT, line=BORDER)
    tf_shape = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.12),
        w - Inches(0.3), h - Inches(0.24))
    tf = tf_shape.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.strip().split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9.5)
        run.font.name = "Courier New"
        run.font.color.rgb = TEXT


def tag(slide, label, x, y, color=ACCENT, w=Inches(1.3)):
    r, g, b = color[0], color[1], color[2]
    bg = RGBColor(min(r + 210, 255), min(g + 210, 255), min(b + 210, 255))
    rect(slide, x, y, w, Inches(0.27), fill=bg, line=color, line_w=Pt(1))
    txt(slide, label, x + Inches(0.06), y + Inches(0.02),
        w - Inches(0.12), Inches(0.23),
        size=9, bold=True, color=color, align=PP_ALIGN.CENTER)


def chip(slide, icon, label, x, y, color=ACCENT):
    r, g, b = color[0], color[1], color[2]
    bg = RGBColor(min(r + 210, 255), min(g + 210, 255), min(b + 210, 255))
    rect(slide, x, y, Inches(2.8), Inches(0.7), fill=bg, line=color, line_w=Pt(1))
    txt(slide, icon + "  " + label,
        x + Inches(0.1), y + Inches(0.14), Inches(2.6), Inches(0.35),
        size=12, bold=True, color=color)


def flow(slide, items, y, bar_color=P1_BAR):
    n     = len(items)
    box_w = Inches(2.4)
    gap   = Inches(0.55)
    total = n * box_w + (n - 1) * gap
    x0    = (W - total) / 2

    for i, (lbl, sub) in enumerate(items):
        bx = x0 + i * (box_w + gap)
        rect(slide, bx, y, box_w, Inches(0.8), fill=PANEL, line=bar_color)
        txt(slide, lbl, bx + Inches(0.1), y + Inches(0.06),
            box_w - Inches(0.2), Inches(0.35),
            size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        if sub:
            txt(slide, sub, bx + Inches(0.1), y + Inches(0.44),
                box_w - Inches(0.2), Inches(0.3),
                size=9, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
        if i < n - 1:
            txt(slide, "→",
                bx + box_w + Inches(0.1), y + Inches(0.22),
                Inches(0.35), Inches(0.35),
                size=16, color=bar_color, align=PP_ALIGN.CENTER)


def stat_row(slide, stats, y, bar_color=P1_BAR):
    n  = len(stats)
    bw = Inches(12.33 / n - 0.15)
    for i, (num, lbl) in enumerate(stats):
        bx = Inches(0.35 + i * (12.33 / n))
        rect(slide, bx, y, bw, Inches(1.0), fill=PANEL, line=BORDER)
        txt(slide, num, bx, y + Inches(0.06), bw, Inches(0.55),
            size=28, bold=True, color=bar_color, align=PP_ALIGN.CENTER)
        txt(slide, lbl, bx, y + Inches(0.6), bw, Inches(0.35),
            size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)


def data_table(slide, headers, rows, x, y, w, h, col_widths=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    t = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = cw
    row_h = h // n_rows
    for ri in range(n_rows):
        t.rows[ri].height = row_h
    for ci, hdr in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = hdr
        run.font.size = Pt(11); run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Segoe UI"
    for ri, row_data in enumerate(rows, 1):
        for ci, val in enumerate(row_data):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_ALT if ri % 2 == 0 else PANEL
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.size = Pt(10); run.font.color.rgb = TEXT_SUB
            run.font.name = "Segoe UI"


def screenshot(slide, filename, x, y, w, h):
    path = os.path.join(SCREENSHOTS, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        r = rect(slide, x, y, w, h, fill=PANEL_ALT, line=BORDER)
        txt(slide, "[ screenshot: " + filename + " ]",
            x, y + h/2 - Inches(0.2), w, Inches(0.4),
            size=9, color=TEXT_MUTE, align=PP_ALIGN.CENTER)


def divider_slide(title, subtitle, bar_color=P1_BAR):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = TEXT
    no_line(rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5),
                 fill=bar_color))
    txt(s, title,
        Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.2),
        size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(s, subtitle,
        Inches(0.6), Inches(3.85), Inches(12.1), Inches(0.6),
        size=18, color=RGBColor(0xc0, 0xc8, 0xd8), align=PP_ALIGN.LEFT)
    txt(s, "GNR629  ·  CSRE, IIT Bombay",
        Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
        size=11, color=RGBColor(0x5c, 0x68, 0x82), align=PP_ALIGN.LEFT)
    return s


# ═════════════════════════════════════════════════════════════════════════════
#  PART 1: PROJECT 1 — INTEROPERABLE GIS SYSTEM (slides 1–15)
# ═════════════════════════════════════════════════════════════════════════════

# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = new_slide()
no_line(rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5), fill=P1_BAR))
no_line(rect(s, Inches(0.3), Inches(0), Inches(13.03), Inches(0.08), fill=BORDER))

txt(s, "PROJECT 1",
    Inches(0.6), Inches(1.2), Inches(11.0), Inches(0.5),
    size=13, bold=True, color=P1_BAR)
txt(s, "Interoperable GIS System",
    Inches(0.6), Inches(1.7), Inches(11.0), Inches(1.3),
    size=44, bold=True, color=TEXT)
txt(s, "An OGC-compliant WMS / WFS web client for geospatial data visualisation",
    Inches(0.6), Inches(3.0), Inches(11.0), Inches(0.6),
    size=17, color=TEXT_SUB)
no_line(rect(s, Inches(0.6), Inches(3.75), Inches(4.0), Inches(0.04), fill=BORDER))
txt(s, "GNR629  ·  CSRE, IIT Bombay  ·  Manas Avinashe  ·  2026",
    Inches(0.6), Inches(3.9), Inches(11.0), Inches(0.4),
    size=12, color=TEXT_MUTE)
for i, (lbl, col) in enumerate([("OGC WMS", ACCENT), ("OGC WFS", GREEN),
                                  ("GeoServer", PURPLE), ("OpenLayers", AMBER)]):
    tag(s, lbl, Inches(0.6 + i * 1.5), Inches(4.6), col)


# ── Slide 2: Problem & Motivation ─────────────────────────────────────────────
s = new_slide()
heading(s, "Problem & Motivation", "Why interoperability in GIS?")

card(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(4.0),
     "The Challenge", [
     "Geospatial data is distributed across",
     " many servers and formats",
     "Different vendors use incompatible APIs",
     "No standard way to discover, request,",
     " or render spatial layers remotely",
     "Manual download + re-processing is slow",
], title_color=RED)

card(s, Inches(6.55), Inches(1.35), Inches(6.43), Inches(4.0),
     "The Solution — OGC Web Services", [
     "WMS — Web Map Service",
     " Render layers as images on demand",
     "WFS — Web Feature Service",
     " Retrieve vector features as data",
     "Standard URL operations across all servers",
     "Any compliant client talks to any server",
], title_color=GREEN)

no_line(rect(s, Inches(0.35), Inches(5.55), Inches(12.63), Inches(0.04), fill=BORDER))
txt(s, "Application Domain: Soil Moisture Monitoring — Maharashtra, India",
    Inches(0.35), Inches(5.7), Inches(12.63), Inches(0.4),
    size=12, bold=True, color=TEXT_SUB)
txt(s, "Soil moisture is critical for agricultural productivity, drought assessment, and irrigation planning."
       " This project serves Maharashtra soil moisture layers through GeoServer and visualises them via the OGC client.",
    Inches(0.35), Inches(6.1), Inches(12.63), Inches(0.65),
    size=11, color=TEXT_MUTE)


# ── Slide 3: Technology Stack ──────────────────────────────────────────────────
s = new_slide()
heading(s, "Technology Stack")

for i, (title, items, col) in enumerate([
    ("Server Side", [
        "Apache Tomcat — application server",
        "GeoServer 2.x — OGC WMS/WFS server",
        "PostgreSQL + PostGIS — spatial database",
        "Stores Maharashtra layers (countries,",
        " mh_villages, soil_moisture_georef)",
    ], ACCENT),
    ("Client Side", [
        "HTML5 / CSS3 / Vanilla JavaScript",
        "OpenLayers — interactive mapping library",
        "DOMParser — built-in XML parsing",
        "No build tools, no npm, no framework",
        "Single index.html opens in any browser",
    ], GREEN),
    ("OGC Standards", [
        "WMS 1.1 — GetCapabilities, GetMap,",
        " GetFeatureInfo",
        "WFS 2.0 — GetCapabilities, GetFeature",
        "EPSG:4326 and EPSG:3857 CRS support",
        "XML response parsing + display",
    ], PURPLE),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(1.35), Inches(4.1), Inches(4.9),
         title, items, title_color=col)

stat_row(s, [
    ("4", "JS modules"), ("2", "OGC standards"), ("3", "GeoServer layers"), ("0", "npm dependencies")
], Inches(6.5))


# ── Slide 4: System Architecture ──────────────────────────────────────────────
s = new_slide()
heading(s, "System Architecture")

flow(s, [
    ("Browser", "index.html\nOpenLayers"),
    ("HTTP/S", "GET request\nURL params"),
    ("GeoServer", "Tomcat:8080\nPostGIS backend"),
    ("Response", "image/png\nor XML"),
    ("Render", "OL tile layer\nor vector"),
], Inches(1.9), P1_BAR)

for i, (title, items, col) in enumerate([
    ("WMS Flow", [
        "GetCapabilities → parse layer list",
        "User selects layer",
        "GetMap → image URL → OL TileWMS",
        "Image rendered as map overlay",
    ], ACCENT),
    ("WFS Flow", [
        "GetCapabilities → parse feature types",
        "GetFeature → GeoJSON response",
        "OL VectorSource + GeoJSON format",
        "Features rendered on map + table",
    ], GREEN),
    ("Click / Info Flow", [
        "Vector click → attributes from OL",
        "WMS click → GetFeatureInfo request",
        "AbortController cancels in-flight reqs",
        "Attributes shown in parsed data panel",
    ], PURPLE),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(3.15), Inches(4.1), Inches(3.5),
         title, items, title_color=col, bullet_size=10)


# ── Slide 5: OGC Standards — WMS & WFS ────────────────────────────────────────
s = new_slide()
heading(s, "OGC Standards — WMS & WFS",
        "How the client communicates with GeoServer")

card(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(5.4),
     "WMS — Web Map Service", [
     "Returns rendered map images (PNG/JPEG)",
     "",
     "Key operations:",
     " GetCapabilities — list available layers",
     " GetMap — request a map image",
     "   LAYERS, BBOX, SRS, WIDTH, HEIGHT",
     " GetFeatureInfo — query features by pixel",
     "",
     "URL example:",
     " ?SERVICE=WMS&REQUEST=GetMap",
     " &LAYERS=gnr629:countries",
     " &BBOX=68,6,98,38&SRS=EPSG:4326",
     " &WIDTH=800&HEIGHT=600&FORMAT=image/png",
], title_color=ACCENT)

card(s, Inches(6.75), Inches(1.35), Inches(6.23), Inches(5.4),
     "WFS — Web Feature Service", [
     "Returns raw vector feature data",
     "",
     "Key operations:",
     " GetCapabilities — list feature types",
     " GetFeature — download features as",
     "   GeoJSON, GML, or KML",
     "   TYPENAMES, MAXFEATURES, BBOX",
     "",
     "URL example:",
     " ?SERVICE=WFS&REQUEST=GetFeature",
     " &TYPENAMES=gnr629:mh_villages",
     " &OUTPUTFORMAT=application/json",
     " &COUNT=50",
], title_color=GREEN)


# ── Slide 6: GetCapabilities ───────────────────────────────────────────────────
s = new_slide()
heading(s, "GetCapabilities — Service Discovery",
        "The first request — answers: what layers does this server have?")

code_card(s, """\
<!-- WMS GetCapabilities response — parsed by xmlParser.js -->

<WMS_Capabilities version="1.1.1">
  <Service>
    <Title>GeoServer Web Map Service</Title>
  </Service>
  <Capability>
    <Layer>
      <Layer queryable="1">
        <Name>gnr629:countries</Name>
        <Title>World Countries</Title>
        <SRS>EPSG:4326</SRS>
        <LatLonBoundingBox minx="-180" miny="-90" maxx="180" maxy="90"/>
      </Layer>
      <Layer queryable="1">
        <Name>gnr629:mh_villages</Name>
        <Title>Maharashtra Villages</Title>
        <LatLonBoundingBox minx="72.6" miny="15.6" maxx="80.9" maxy="22.0"/>
      </Layer>
    </Layer>
  </Capability>
</WMS_Capabilities>""",
    Inches(0.35), Inches(1.35), Inches(6.8), Inches(5.4))

card(s, Inches(7.35), Inches(1.35), Inches(5.63), Inches(2.6),
     "What the client does with this", [
     "DOMParser parses the raw XML",
     "Extracts layer Name, Title, SRS, BBOX",
     "Populates the layer dropdown",
     "Auto-fills BBOX fields on layer select",
     "Auto-fills SRS field from layer CRS",
], title_color=ACCENT)

card(s, Inches(7.35), Inches(4.15), Inches(5.63), Inches(2.6),
     "Key JS — xmlParser.js", [
     "doc = new DOMParser()",
     "      .parseFromString(xml, 'text/xml')",
     "layers = doc.getElementsByTagName('Layer')",
     "name   = el.querySelector('Name').textContent",
     "bbox   = el.querySelector('[minx]')",
], title_color=PURPLE)


# ── Slide 7: WMS GetMap ────────────────────────────────────────────────────────
s = new_slide()
heading(s, "WMS GetMap — Map Visualisation",
        "Render any server layer as an image overlay on the base map")

screenshot(s, "8. getMap_countries.png",
           Inches(0.35), Inches(1.35), Inches(5.5), Inches(4.3))

card(s, Inches(6.1), Inches(1.35), Inches(6.88), Inches(2.1),
     "GetMap Request Built by ogcRequests.js", [
     "SERVICE=WMS & REQUEST=GetMap",
     "LAYERS=<selected layer>",
     "BBOX=minx,miny,maxx,maxy (EPSG:4326)",
     "SRS=EPSG:4326  WIDTH=800  HEIGHT=600",
     "FORMAT=image/png (or jpeg/tiff)",
], title_color=ACCENT)

card(s, Inches(6.1), Inches(3.65), Inches(6.88), Inches(2.05),
     "OpenLayers Integration", [
     "ol.layer.Image + ol.source.ImageWMS",
     "Layer added to map with auto z-index",
     "Multiple layers stack correctly",
     "Layer checklist panel tracks all layers",
], title_color=GREEN)

screenshot(s, "9. getMap_mh_villages.png",
           Inches(6.1), Inches(5.85), Inches(6.88), Inches(1.4))


# ── Slide 8: WFS GetFeature ────────────────────────────────────────────────────
s = new_slide()
heading(s, "WFS GetFeature — Vector Data",
        "Retrieve actual feature geometry and attributes, not just a rendered image")

screenshot(s, "11. getFeature_countries_258.png",
           Inches(0.35), Inches(1.35), Inches(5.5), Inches(4.3))

card(s, Inches(6.1), Inches(1.35), Inches(6.88), Inches(2.1),
     "GetFeature Request", [
     "SERVICE=WFS & REQUEST=GetFeature",
     "TYPENAMES=<selected layer>",
     "OUTPUTFORMAT=application/json",
     "COUNT=<maxFeatures param>",
     "BBOX=<if set by user>",
], title_color=GREEN)

card(s, Inches(6.1), Inches(3.65), Inches(6.88), Inches(2.05),
     "Result", [
     "GeoJSON features loaded into OL VectorLayer",
     "Features rendered as polygons/points",
     "Attributes displayed in parsed data table",
     "Click a feature → highlight + show info",
], title_color=ACCENT)

screenshot(s, "12. getFeature_countries_50.png",
           Inches(6.1), Inches(5.85), Inches(6.88), Inches(1.4))


# ── Slide 9: GetFeatureInfo / Click Query ─────────────────────────────────────
s = new_slide()
heading(s, "GetFeatureInfo — Click to Query",
        "Click anywhere on the map → retrieve feature attributes at that pixel")

card(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.4),
     "How It Works", [
     "User clicks on the map",
     "",
     "1. If a vector (WFS) feature is at that",
     "   pixel → attributes shown immediately",
     "   from the in-memory GeoJSON",
     "",
     "2. If no vector feature → GetFeatureInfo",
     "   request sent to WMS server:",
     "   REQUEST=GetFeatureInfo",
     "   QUERY_LAYERS=<layer>",
     "   I=<pixel X>  J=<pixel Y>",
     "",
     "AbortController cancels any previous",
     "in-flight request on each new click",
], title_color=ACCENT)

code_card(s, """\
// map.js — WMS click handler

sosMap.on("singleclick", async (evt) => {
  const viewRes = view.getResolution();
  const url = wmsSource.getFeatureInfoUrl(
    evt.coordinate, viewRes, "EPSG:4326",
    { INFO_FORMAT: "application/json",
      FEATURE_COUNT: maxFeatures }
  );

  // cancel previous in-flight request
  if (abortCtrl) abortCtrl.abort();
  abortCtrl = new AbortController();

  const res = await fetch(url, { signal: abortCtrl.signal });
  const json = await res.json();
  // display features in parsed data panel
  displayFeatureInfo(json.features);
});""",
    Inches(6.5), Inches(1.35), Inches(6.48), Inches(5.4))


# ── Slide 10: BBOX Selection ──────────────────────────────────────────────────
s = new_slide()
heading(s, "Bounding Box — Three Ways to Set AOI",
        "Area of interest controls both the map view and the request parameters")

for i, (title, bullets, col, sc) in enumerate([
    ("Method 1: Auto-fill on Layer Select",
     ["Click a layer in the dropdown",
      "BBOX auto-fills from GetCapabilities",
      " LatLonBoundingBox parsed by xmlParser.js",
      "SRS also auto-populated",
     ], ACCENT, "19. bbox_selection.png"),
    ("Method 2: Draw on Map",
     ["Click '✏ Draw on Map' button",
      "Drag rectangle on OpenLayers map",
      "ol.interaction.Draw + createBox()",
      "Coords converted EPSG:3857 → EPSG:4326",
     ], GREEN, "20. bbox_update.png"),
    ("Method 3: Manual + Apply",
     ["Type Min Lon/Lat + Max Lon/Lat directly",
      "Click 'Apply to Map'",
      "Validates ±180° / ±90° ranges",
      "Map view pans and zooms to BBOX",
     ], PURPLE, "21. bbox_apply_to_map.png"),
]):
    bx = Inches(0.35 + i * 4.35)
    card(s, bx, Inches(1.35), Inches(4.1), Inches(2.8), title, bullets, title_color=col)
    screenshot(s, sc, bx, Inches(4.35), Inches(4.1), Inches(2.8))


# ── Slide 11: Layer Management ────────────────────────────────────────────────
s = new_slide()
heading(s, "Layer Management",
        "Stack and control multiple WMS/WFS layers simultaneously")

screenshot(s, "10. soil_moisture_layer_stacked.png",
           Inches(0.35), Inches(1.35), Inches(5.5), Inches(5.4))

card(s, Inches(6.1), Inches(1.35), Inches(6.88), Inches(1.9),
     "Added Layers Panel", [
     "Lists every loaded WMS / WFS layer",
     "WMS / WFS badge shown per layer",
     "Toggle visibility on/off",
     "Adjust opacity with slider",
], title_color=ACCENT)

card(s, Inches(6.1), Inches(3.45), Inches(6.88), Inches(1.9),
     "Z-Index & Stacking", [
     "Most recently added layer on top",
     "WMS layers use auto-incrementing zIndex",
     "WFS vector layers always above WMS images",
     "Clear All Layers button resets map",
], title_color=GREEN)

card(s, Inches(6.1), Inches(5.55), Inches(6.88), Inches(1.2),
     "Tested Layers", [
     "gnr629:countries   ·   gnr629:mh_villages",
     "gnr629:soil_moisture_georef",
], title_color=PURPLE)


# ── Slide 12: XML Parsing ─────────────────────────────────────────────────────
s = new_slide()
heading(s, "XML Parsing — xmlParser.js",
        "DOMParser turns raw XML responses into structured data")

code_card(s, """\
// xmlParser.js — parse WMS GetCapabilities

function parseWMSCapabilities(xmlText) {
  const doc = new DOMParser()
               .parseFromString(xmlText, "text/xml");

  const layerEls = doc.querySelectorAll("Layer > Layer");
  const layers   = [];

  layerEls.forEach(el => {
    const name  = el.querySelector("Name")?.textContent;
    const title = el.querySelector("Title")?.textContent;
    const bbox  = el.querySelector("LatLonBoundingBox");
    const srs   = el.querySelector("SRS")?.textContent;

    if (!name) return;

    layers.push({
      name,  title,  srs,
      minx: parseFloat(bbox?.getAttribute("minx")),
      miny: parseFloat(bbox?.getAttribute("miny")),
      maxx: parseFloat(bbox?.getAttribute("maxx")),
      maxy: parseFloat(bbox?.getAttribute("maxy")),
    });
  });

  return layers;   // → populates layer dropdown
}

// Copy raw XML to response panel
function displayXMLResponse(xmlText) {
  document.getElementById("xmlOutput").textContent = xmlText;
}""",
    Inches(0.35), Inches(1.35), Inches(7.0), Inches(5.6))

card(s, Inches(7.55), Inches(1.35), Inches(5.43), Inches(2.6),
     "DOMParser API", [
     "Built into every browser — no libraries",
     "parseFromString(str, 'text/xml')",
     "Returns a navigable DOM tree",
     "querySelector() works on XML too",
     "textContent gives element text value",
], title_color=ACCENT)

card(s, Inches(7.55), Inches(4.15), Inches(5.43), Inches(2.8),
     "Parsed Data Table", [
     "WFS attributes rendered in a 2-col table",
     "Property | Value rows",
     "WMS parsed layers → layer checklist",
     "GetFeatureInfo JSON → attribute rows",
     "Raw XML always shown in xml panel",
], title_color=GREEN)


# ── Slide 13: UI Features ─────────────────────────────────────────────────────
s = new_slide()
heading(s, "User Interface",
        "Clean, professional design with light/dark themes")

screenshot(s, "4. client_whiteTheme.png",
           Inches(0.35), Inches(1.35), Inches(6.0), Inches(3.5))
screenshot(s, "5. client_darkTheme.png",
           Inches(6.65), Inches(1.35), Inches(6.33), Inches(3.5))

txt(s, "Light Theme", Inches(0.35), Inches(4.9), Inches(6.0), Inches(0.3),
    size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
txt(s, "Dark Theme", Inches(6.65), Inches(4.9), Inches(6.33), Inches(0.3),
    size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)

for i, feat in enumerate([
    "Light / dark theme toggle — saved in localStorage",
    "Resizable map / response panes — drag divider",
    "Collapsible map & response panels",
    "Status dot: idle / loading / active / error",
    "Toast notification system",
]):
    txt(s, "• " + feat,
        Inches(0.35), Inches(5.25 + i * 0.32), Inches(12.63), Inches(0.3),
        size=11, color=TEXT_SUB)


# ── Slide 14: GeoServer Workspace & External Servers ─────────────────────────
s = new_slide()
heading(s, "GeoServer Workspace & External Servers",
        "Local PostGIS layers and remote public WMS services")

screenshot(s, "0. gnr629_workspace.png",
           Inches(0.35), Inches(1.35), Inches(4.2), Inches(3.2))

card(s, Inches(4.75), Inches(1.35), Inches(8.23), Inches(3.2),
     "Local GeoServer (localhost:8080)", [
     "Workspace: gnr629",
     "Layers stored in PostgreSQL/PostGIS:",
     " countries — world boundaries",
     " mh_villages — Maharashtra village polygons",
     " soil_moisture_georef — raster soil data",
     "Served via WMS and WFS endpoints",
], title_color=ACCENT)

screenshot(s, "16. external_wms_server_nasa.png",
           Inches(0.35), Inches(4.75), Inches(4.2), Inches(2.5))

card(s, Inches(4.75), Inches(4.75), Inches(8.23), Inches(2.5),
     "External Public WMS Servers (tested)", [
     "NASA GIBS — gibs.earthdata.nasa.gov",
     "NOAA Weather — opengeo.ncep.noaa.gov",
     "GeoSolutions — gs-stable.geo-solutions.it",
     "Client works with any OGC-compliant server",
], title_color=GREEN)


# ── Slide 15: Project 1 Summary ────────────────────────────────────────────────
s = new_slide()
heading(s, "Project 1 — Summary & Results")

stat_row(s, [
    ("3", "GeoServer layers"),
    ("4", "OGC operations"),
    ("3", "External WMS tested"),
    ("2", "Themes"),
    ("0", "npm dependencies"),
], Inches(1.35))

card(s, Inches(0.35), Inches(2.7), Inches(5.9), Inches(3.9),
     "Requirements Met", [
     "✓ WMS GetCapabilities, GetMap, GetFeatureInfo",
     "✓ WFS GetCapabilities, GetFeature",
     "✓ XML response displayed raw + parsed",
     "✓ Multiple stacked layers with z-index",
     "✓ OSM base layer",
     "✓ External WMS/WFS server support",
     "✓ BBOX coordinate input + draw + apply",
     "✓ Format and size selection",
     "✓ Light/dark theme + resizable panels",
], title_color=GREEN)

card(s, Inches(6.55), Inches(2.7), Inches(6.43), Inches(3.9),
     "Foundation for Project 2", [
     "DOMParser replaces any XML library",
     "OpenLayers handles CRS transforms internally",
     "AbortController prevents stale responses",
     "CSS variables make theme switching trivial",
     "",
     "Same map library, same UI patterns",
     "Same XML parsing approach",
     "Extended with SOS tab + Docker + PostGIS",
], title_color=ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
#  PART 2: PROJECT 2 — GeoPulse istSOS + Docker + PostGIS (slides 16–30)
# ═════════════════════════════════════════════════════════════════════════════

# ── Slide 16: Title — Project 2 ───────────────────────────────────────────────
s = new_slide()
no_line(rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5), fill=P2_BAR))
no_line(rect(s, Inches(0.3), Inches(0), Inches(13.03), Inches(0.08), fill=BORDER))

txt(s, "PROJECT 2",
    Inches(0.6), Inches(1.2), Inches(11.0), Inches(0.5),
    size=13, bold=True, color=P2_BAR)
txt(s, "GeoPulse",
    Inches(0.6), Inches(1.7), Inches(11.0), Inches(1.3),
    size=44, bold=True, color=TEXT)
txt(s, "OGC Sensor Observation Service — istSOS 3 · Docker Compose · PostGIS",
    Inches(0.6), Inches(3.0), Inches(11.0), Inches(0.6),
    size=17, color=TEXT_SUB)
no_line(rect(s, Inches(0.6), Inches(3.75), Inches(4.0), Inches(0.04), fill=BORDER))
txt(s, "GNR629  ·  CSRE, IIT Bombay  ·  Manas Avinashe  ·  2026",
    Inches(0.6), Inches(3.9), Inches(11.0), Inches(0.4),
    size=12, color=TEXT_MUTE)
for i, (lbl, col) in enumerate([
    ("istSOS 3", TEAL), ("Docker", DOCKER), ("PostGIS", PG_BLUE),
    ("OGC SOS 2.0", GREEN), ("SensorML", ACCENT)
]):
    tag(s, lbl, Inches(0.6 + i * 1.6), Inches(4.6), col, w=Inches(1.45))


# ── Slide 17: Why istSOS & Docker? ────────────────────────────────────────────
s = new_slide()
heading(s, "Why istSOS + Docker + PostGIS?",
        "Upgrading from a custom FastAPI/SQLite prototype to a production-grade OGC stack",
        bar_color=P2_BAR)

flow(s, [
    ("FastAPI\n+ SQLite", "Custom SOS API\nflat file DB"),
    ("Limitations", "No spatial index\nmanual XML\nfragile paths"),
    ("Upgrade To", "Real OGC server\nSpatial DB\nContainerised"),
    ("istSOS 3\n+ PostGIS", "Standards-native\nST_* functions\nDocker Compose"),
], Inches(1.7), P2_BAR)

card(s, Inches(0.35), Inches(3.1), Inches(5.9), Inches(3.5),
     "Custom FastAPI Approach (v1)", [
     "Hand-written XML templates",
     "SQLite — no native spatial support",
     "BBOX filtering: Python string compare",
     "Path bug: relative CSV paths break",
     "One-file deployment, hard to extend",
     "No real GetCapabilities negotiation",
], title_color=RED)

card(s, Inches(6.55), Inches(3.1), Inches(6.43), Inches(3.5),
     "istSOS + PostGIS Approach (v2)", [
     "istSOS 3 generates OGC-native XML",
     "PostGIS geometry column (GIST index)",
     "BBOX: ST_MakeEnvelope — true spatial query",
     "Docker Compose — reproducible environment",
     "3 services: postgis · backend · loader",
     "OGC URN identifiers for all properties",
], title_color=P2_BAR)


# ── Slide 18: Docker Compose Architecture ─────────────────────────────────────
s = new_slide()
heading(s, "Docker Compose — Three-Service Architecture",
        "docker compose up  →  everything starts, data loads, API is live",
        bar_color=P2_BAR)

flow(s, [
    ("postgis", "postgres:17\nPostGIS 3.5\nport 5432"),
    ("loader", "one-shot\nCSV → PostGIS\nthen exits"),
    ("backend", "FastAPI\nport 8002\nproxies istSOS"),
    ("istSOS", "external\nport 8088\nOGC SOS 2.0"),
    ("Browser", "port 8002\n/app/index.html"),
], Inches(1.8), DOCKER)

code_card(s, """\
# docker-compose.yml — simplified

services:
  postgis:
    image: postgis/postgis:17-3.5
    volumes:
      - postgis-data:/var/lib/postgresql/data
      - ./sql/01-init-postgis.sql:/docker-entrypoint-initdb.d/init.sql
    healthcheck:
      test: pg_isready -U postgres
      interval: 10s

  backend:
    build: ./backend
    ports: ["8002:8000"]
    depends_on:
      postgis: { condition: service_healthy }
    environment:
      - ISTSOS_SERVICE_URL=http://host.docker.internal:8088/istsos/geopulse

  loader:
    build: ./backend
    command: python /app/scripts/load_postgis.py --csv /app/data/cleaned_weather.csv
    depends_on:
      postgis: { condition: service_healthy }
    restart: "no"   # ← runs once, then done""",
    Inches(0.35), Inches(3.1), Inches(12.63), Inches(4.15))


# ── Slide 19: PostGIS Schema & Spatial Indexing ───────────────────────────────
s = new_slide()
heading(s, "PostGIS Schema — Spatial Database Design",
        "geometry(Point, 4326) column enables true spatial queries via GIST index",
        bar_color=P2_BAR)

code_card(s, """\
-- sql/01-init-postgis.sql

CREATE SCHEMA geopulse;

CREATE TABLE geopulse.observations (
  id           BIGSERIAL PRIMARY KEY,
  sensor_id    TEXT NOT NULL,          -- e.g. "Tokyo_Japan"
  location_name TEXT,
  country      TEXT,
  observed_at  TIMESTAMPTZ NOT NULL,   -- UTC timestamp with zone
  temperature  DOUBLE PRECISION,       -- °C
  humidity     INTEGER,                -- %
  wind_speed   DOUBLE PRECISION,       -- kph
  pressure     DOUBLE PRECISION,       -- mb
  precipitation DOUBLE PRECISION,      -- mm
  geom         geometry(Point, 4326) NOT NULL   -- WGS84 point
);

-- Indexes for fast filtering
CREATE INDEX ON geopulse.observations (sensor_id);
CREATE INDEX ON geopulse.observations (country);
CREATE INDEX ON geopulse.observations (observed_at);
CREATE INDEX USING GIST ON geopulse.observations (geom);
--            ^^^^^^^^^^  spatial index — makes ST_MakeEnvelope fast""",
    Inches(0.35), Inches(1.35), Inches(7.7), Inches(5.6))

card(s, Inches(8.1), Inches(1.35), Inches(4.88), Inches(2.6),
     "Why PostGIS vs SQLite", [
     "geometry column stores lon/lat as WKT",
     "GIST index — R-tree spatial structure",
     "  O(log n) bbox intersection queries",
     "ST_MakeEnvelope — native bbox object",
     "geom && envelope — index scan, not full-table",
     "ST_SetSRID(ST_MakePoint(lon,lat), 4326)",
     "  inserts coordinates correctly",
], title_color=PG_BLUE)

card(s, Inches(8.1), Inches(4.15), Inches(4.88), Inches(2.8),
     "Data Loaded By loader Service", [
     "psycopg2.extras.execute_values()",
     "Batch insert: page_size=1000",
     "TRUNCATE RESTART IDENTITY first",
     "17,775 rows across 25 sensors",
     "ST_SetSRID(ST_MakePoint(lon,lat),4326)",
     "  per-row in the INSERT",
], title_color=TEAL)


# ── Slide 20: Data Pipeline ────────────────────────────────────────────────────
s = new_slide()
heading(s, "Data Pipeline — CSV to PostGIS",
        "One command loads 17,775 weather observations into a spatially indexed database",
        bar_color=P2_BAR)

flow(s, [
    ("GlobalWeatherRepository", "138 k rows\n257 cities · 35 MB"),
    ("process_data.py", "select 25 cities\nkeep full history"),
    ("cleaned_weather.csv", "17,775 rows\n2 MB"),
    ("loader container", "psycopg2 batch\ninsert + geom"),
    ("geopulse.observations", "PostGIS\nGIST indexed"),
], Inches(1.7), TEAL)

code_card(s, """\
# scripts/load_postgis.py — core insert logic

import psycopg2, psycopg2.extras, csv, argparse

conn = psycopg2.connect(**pg_kwargs)
cur  = conn.cursor()

cur.execute("TRUNCATE TABLE geopulse.observations RESTART IDENTITY")

INSERT_SQL = \"\"\"
INSERT INTO geopulse.observations
  (sensor_id, location_name, country, observed_at,
   temperature, humidity, wind_speed, pressure, precipitation, geom)
VALUES %s
\"\"\"

rows = []
with open(args.csv) as f:
    for r in csv.DictReader(f):
        rows.append((
            r["sensor_id"], r["location_name"], r["country"], r["timestamp"],
            r["temperature"], r["humidity"], r["wind_speed"],
            r["pressure"], r["precipitation"],
            f"SRID=4326;POINT({r['longitude']} {r['latitude']})"
        ))

psycopg2.extras.execute_values(cur, INSERT_SQL, rows, page_size=1000)
conn.commit()
print(f"Loaded {len(rows)} rows")""",
    Inches(0.35), Inches(3.05), Inches(12.63), Inches(4.15))


# ── Slide 21: OGC Standards Stack ─────────────────────────────────────────────
s = new_slide()
heading(s, "OGC Standards Stack — Project 2",
        "istSOS implements SOS 2.0 natively; our FastAPI backend bridges PostGIS to the same interface",
        bar_color=P2_BAR)

for i, (title, items, col) in enumerate([
    ("SOS 2.0 — Sensor Observation Service", [
     "Defines HOW to request sensor data",
     "Three core operations:",
     " GetCapabilities — service metadata",
     " DescribeSensor  — sensor identity",
     " GetObservation  — observation data",
     "Spatial, temporal, property filters",
     "Implemented natively by istSOS 3",
    ], TEAL),
    ("O&M — Observations & Measurements", [
     "Defines WHAT a sensor reading is:",
     " featureOfInterest — WHERE (GML Point)",
     " observedProperty  — WHAT (OGC URN)",
     " procedure         — WHICH sensor",
     " phenomenonTime    — WHEN (ISO 8601)",
     " result            — THE VALUE + unit",
     "ISO 19156 standard",
    ], GREEN),
    ("SensorML 2.0 + istSOS URNs", [
     "Sensor described in SensorML 2.0 XML",
     "OGC URN identifiers per property:",
     " urn:ogc:def:parameter:x-istsos:",
     "   1.0:meteo:air:temperature",
     "   1.0:meteo:air:humidity",
     "   1.0:meteo:air:wind_speed",
     "Registered via InsertSensor in istSOS",
    ], ACCENT),
]):
    card(s, Inches(0.35 + i * 4.35), Inches(1.35), Inches(4.1), Inches(5.6),
         title, items, title_color=col)


# ── Slide 22: FastAPI Backend Architecture ────────────────────────────────────
s = new_slide()
heading(s, "FastAPI Backend — API Endpoints",
        "Python server bridges the browser to PostGIS (observations) and istSOS (OGC operations)",
        bar_color=P2_BAR)

card(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(5.6),
     "API Endpoints", [
     "GET /",
     "  → redirect to /app/index.html",
     "",
     "GET /health",
     "  → { status, row_count, istsos_url }",
     "",
     "GET /api/config",
     "  → { istsosServiceUrl, defaultProcedure }",
     "",
     "GET /api/sensors",
     "  → 25 sensor locations (JSON)",
     "",
     "GET /api/observations",
     "  → filtered weather data (JSON or XML)",
     "",
     "GET /api/istsos/capabilities",
     "GET /api/istsos/describe-sensor",
     "GET /api/istsos/raw",
     "  → proxy to external istSOS server",
     "",
     "GET /app/*",
     "  → static frontend files",
], title_color=TEAL, bullet_size=10)

card(s, Inches(6.55), Inches(1.35), Inches(6.43), Inches(2.6),
     "psycopg2 Connection", [
     "psycopg2.connect(host, dbname, user, ...)",
     "conn.cursor(cursor_factory=RealDictCursor)",
     "  → rows as dict, not tuple",
     "Env vars: POSTGRES_HOST / _DB / _USER",
     "All queries use %s placeholders",
     "  — SQL injection safe",
], title_color=PG_BLUE)

card(s, Inches(6.55), Inches(4.15), Inches(6.43), Inches(2.8),
     "istSOS Proxy (httpx)", [
     "httpx.get(istsos_url, params=...) ",
     "Forwards SOS requests to external server",
     "Returns raw XML — browser parses it",
     "ISTSOS_SERVICE_URL from .env",
     "Falls back gracefully if istSOS offline",
     "Allows future switch to hosted istSOS",
], title_color=ACCENT)


# ── Slide 23: Spatial Query Engine ────────────────────────────────────────────
s = new_slide()
heading(s, "Spatial Query Engine — ST_MakeEnvelope",
        "PostGIS BBOX filter uses the GIST index — sub-millisecond on 17 k rows",
        bar_color=P2_BAR)

code_card(s, """\
# backend/app/main.py — /api/observations BBOX filter

@app.get("/api/observations")
async def get_observations(
    bbox:    str  = Query(None),   # "minLon,minLat,maxLon,maxLat"
    country: str  = Query(None),
    start:   str  = Query(None),
    end:     str  = Query(None),
    param:   str  = Query(None),   # temperature|humidity|...
    op:      str  = Query(None),   # eq|lt|gt|between|...
    value:   float= Query(None),
    value2:  float= Query(None),
    limit:   int  = Query(500),
    fmt:     str  = Query("json"),
):
    clauses, args = [], []

    if bbox:
        min_lon, min_lat, max_lon, max_lat = map(float, bbox.split(","))
        clauses.append(
            "geom && ST_MakeEnvelope(%s, %s, %s, %s, 4326)"
        )
        args.extend([min_lon, min_lat, max_lon, max_lat])
        # ↑ && operator triggers GIST index — no full table scan

    if country:
        clauses.append("LOWER(country) = LOWER(%s)")
        args.append(country)

    if start:
        clauses.append("observed_at >= %s"); args.append(start)
    if end:
        clauses.append("observed_at <= %s"); args.append(end)

    if param and op and value is not None:
        op_map = {"eq":"=","neq":"!=","lt":"<","lte":"<=","gt":">","gte":">="}
        if op == "between":
            clauses.append(f"{param} BETWEEN %s AND %s")
            args.extend([value, value2])
        else:
            clauses.append(f"{param} {op_map[op]} %s"); args.append(value)

    where = "WHERE " + " AND ".join(clauses) if clauses else ""
    sql   = f"SELECT ... FROM geopulse.observations {where} LIMIT %s"
    args.append(limit)""",
    Inches(0.35), Inches(1.35), Inches(12.63), Inches(5.9))


# ── Slide 24: istSOS Integration & OGC Proxy ─────────────────────────────────
s = new_slide()
heading(s, "istSOS Integration — OGC Proxy Endpoints",
        "The backend transparently forwards GetCapabilities / DescribeSensor to a real istSOS server",
        bar_color=P2_BAR)

code_card(s, """\
# backend/app/main.py — istSOS proxy routes

import httpx

ISTSOS_URL = os.getenv("ISTSOS_SERVICE_URL",
    "http://host.docker.internal:8088/istsos/geopulse")

@app.get("/api/istsos/capabilities")
async def istsos_capabilities():
    r = httpx.get(ISTSOS_URL, params={
        "service": "SOS",
        "version": "2.0.0",
        "request": "GetCapabilities"
    })
    return Response(content=r.content,
                    media_type="application/xml")

@app.get("/api/istsos/describe-sensor")
async def istsos_describe_sensor(procedure: str = Query(...)):
    r = httpx.get(ISTSOS_URL, params={
        "service": "SOS", "version": "2.0.0",
        "request": "DescribeSensor",
        "procedure": procedure,
        "procedureDescriptionFormat":
            "http://www.opengis.net/sensorML/1.0.1"
    })
    return Response(content=r.content,
                    media_type="application/xml")

@app.get("/api/istsos/raw")
async def istsos_raw(
    request: str, offering: str = None,
    procedure: str = None, observedProperty: str = None,
    eventTime: str = None, responseFormat: str = None
):
    params = {"service":"SOS","version":"2.0.0","request":request}
    # ... forward any provided params
    r = httpx.get(ISTSOS_URL, params=params)
    return Response(content=r.content, media_type="application/xml")""",
    Inches(0.35), Inches(1.35), Inches(7.5), Inches(5.9))

card(s, Inches(7.85), Inches(1.35), Inches(5.13), Inches(2.6),
     "istSOS 3 — What It Provides", [
     "Full OGC SOS 2.0.0 implementation",
     "InsertSensor — register a procedure",
     "InsertObservation — ingest readings",
     "GetObservation — OGC-native XML output",
     "Admin REST API for service config",
     "Docker image: istsos/istsos3",
], title_color=TEAL)

card(s, Inches(7.85), Inches(4.15), Inches(5.13), Inches(2.8),
     "Import Artifacts (prepared)", [
     "scripts/prepare_istsos_import.py",
     "Generates artifacts/istsos-import/",
     " procedures/ — 25 JSON sensor defs",
     " csv/ — 25 .dat files (istSOS format)",
     "Header uses OGC URN identifiers:",
     " urn:ogc:def:parameter:x-istsos:",
     "   1.0:meteo:air:temperature",
], title_color=ACCENT)


# ── Slide 25: istSOS Import Artifacts ─────────────────────────────────────────
s = new_slide()
heading(s, "istSOS Import Artifacts — OGC Data Format",
        "prepare_istsos_import.py generates sensor definitions and .dat files ready for istSOS ingestion",
        bar_color=P2_BAR)

code_card(s, """\
# artifacts/istsos-import/procedures/Tokyo_Japan.json
{
  "sensor_id": "Tokyo_Japan",
  "location_name": "Tokyo",
  "country": "Japan",
  "latitude": 35.6895,
  "longitude": 139.6917,
  "observed_properties": [
    "urn:ogc:def:parameter:x-istsos:1.0:meteo:air:temperature",
    "urn:ogc:def:parameter:x-istsos:1.0:meteo:air:humidity",
    "urn:ogc:def:parameter:x-istsos:1.0:meteo:air:wind_speed",
    "urn:ogc:def:parameter:x-istsos:1.0:meteo:air:pressure",
    "urn:ogc:def:parameter:x-istsos:1.0:meteo:air:precipitation"
  ]
}

# artifacts/istsos-import/csv/Tokyo_Japan.dat  (first 2 lines)
urn:ogc:def:parameter:x-istsos:1.0:time:iso8601,urn:ogc:def:parameter:x-istsos:1.0:meteo:air:temperature,...
2024-05-16T08:45:00+00:00,25.3,68,4.5,1013.2,0.0""",
    Inches(0.35), Inches(1.35), Inches(7.7), Inches(4.0))

card(s, Inches(0.35), Inches(5.55), Inches(7.7), Inches(1.65),
     "Loading into istSOS", [
     "1. POST /wa/istsos/services/geopulse/procedures — register sensor",
     "2. POST /wa/istsos/services/geopulse/observations — upload .dat file",
     "3. istSOS stores in its own PostgreSQL schema",
     "4. GetObservation returns fully OGC-compliant XML",
], title_color=TEAL, bullet_size=10)

card(s, Inches(8.1), Inches(1.35), Inches(4.88), Inches(5.85),
     "URN Identifier Structure", [
     "urn:ogc:def:parameter",
     "  :x-istsos",
     "  :1.0",
     "  :meteo:air:<property>",
     "",
     "Properties:",
     " :temperature  (°C)",
     " :humidity     (%)",
     " :wind_speed   (kph)",
     " :pressure     (mb)",
     " :precipitation (mm)",
     "",
     "Time:",
     " :time:iso8601",
     "",
     "Same URNs used in:",
     " SensorML outputs",
     " GetObservation filter",
     " .dat file header",
], title_color=ACCENT, bullet_size=10)


# ── Slide 26: Frontend SOS Tab ────────────────────────────────────────────────
s = new_slide()
heading(s, "Frontend — SOS Tab Architecture",
        "Two-panel layout: filter panel on the left, map + table + charts on the right",
        bar_color=P2_BAR)

card(s, Inches(0.35), Inches(1.35), Inches(4.1), Inches(5.6),
     "Left: Filter Panel (300 px)", [
     "Quick Requests row:",
     " istSOS Capabilities  |  DescribeSensor",
     "",
     "Spatial Subsetting:",
     " Country input  (datalist: 25 nations)",
     " BBOX fields  +  Draw on Map button",
     "",
     "Temporal Subsetting:",
     " After datetime OR Start–End window",
     "",
     "Parameter Filter:",
     " Property: temperature / humidity / ...",
     " Operator: eq / lt / gt / between / ...",
     " Value (+ Value2 for between)",
     "",
     "Output: Max Observations (1–5000)",
     "",
     "Query PostGIS button → fetch()",
     "Last Request URL displayed",
     "Raw XML response panel (copy button)",
], title_color=TEAL, bullet_size=10)

card(s, Inches(4.7), Inches(1.35), Inches(4.0), Inches(5.6),
     "Right: Map Pane", [
     "OpenLayers map (OSM base)",
     "25 sensor markers on load",
     " — coloured by avg temperature",
     "Temperature colour scale:",
     " < 0°C  = blue",
     " 0–15   = cyan-green",
     " 15–25  = green",
     " 25–35  = amber",
     " > 35   = red",
     "Click marker → popup with",
     " latest 5 readings",
     "Row click → pan to sensor",
     " marker pulses",
     "BBOX draw tool synced with",
     " filter panel inputs",
], title_color=ACCENT, bullet_size=10)

card(s, Inches(9.0), Inches(1.35), Inches(3.98), Inches(5.6),
     "Right: Table + Charts", [
     "Observation table:",
     " #, City, Country, Timestamp",
     " Temperature badge (color-coded)",
     " Humidity, Wind, Pressure, Precip",
     " Clickable rows",
     "",
     "Chart.js bar chart:",
     " Average temperature per city",
     "",
     "Chart.js doughnut chart:",
     " Observation count per country",
     "",
     "Destroyed + re-rendered on",
     " each new query result",
     "maintainAspectRatio: false",
], title_color=GREEN, bullet_size=10)


# ── Slide 27: sos.js — Key Implementation ─────────────────────────────────────
s = new_slide()
heading(s, "sos.js — Core Implementation",
        "450+ lines managing map state, sensor markers, observations, table, and charts",
        bar_color=P2_BAR)

code_card(s, """\
// sos.js — fetch observations from PostGIS backend

async function fetchObservations() {
  const params = new URLSearchParams();

  const country = document.getElementById("sosCountry").value.trim();
  if (country) params.append("country", country);

  const bboxInputs = ["bboxMinLon","bboxMinLat","bboxMaxLon","bboxMaxLat"];
  const bbox = bboxInputs.map(id => document.getElementById(id).value).filter(Boolean);
  if (bbox.length === 4) params.append("bbox", bbox.join(","));

  const after = document.getElementById("sosAfter").value;
  if (after) params.append("after", after.replace("T"," "));

  const param = document.getElementById("sosParam").value;
  const op    = document.getElementById("sosOp").value;
  const val   = document.getElementById("sosValue").value;
  if (param && op && val) {
    params.append("param", param);
    params.append("op", op);
    params.append("value", val);
  }

  const limit = document.getElementById("sosLimit").value || 500;
  params.append("limit", limit);
  params.append("fmt", "json");    // ← JSON for table/chart rendering

  const url = `${API_BASE}/api/observations?${params}`;
  document.getElementById("lastRequestURL").textContent = url;

  const response = await fetch(url);
  const obs      = await response.json();

  _renderObsMarkers(obs);   // ← update map markers
  _renderTable(obs);         // ← populate table
  _renderCharts(obs);        // ← Chart.js bar + doughnut
}""",
    Inches(0.35), Inches(1.35), Inches(7.5), Inches(5.9))

card(s, Inches(8.0), Inches(1.35), Inches(4.98), Inches(5.9),
     "State Variables", [
     "sosMap — OL map instance",
     "sosMarkerSource — OL VectorSource",
     "  holds all 25 sensor features",
     "_currentObs — last observation array",
     "_sensorMeta — dict of sensor metadata",
     "  keyed by sensor_id",
     "_barChart — Chart.js instance",
     "_pieChart — Chart.js instance",
     "",
     "API_BASE =",
     " '' if served by FastAPI",
     " 'http://127.0.0.1:8002'",
     " if opened as file://",
     "",
     "_renderObsMarkers() sets marker",
     " style based on avg temperature",
     " uses temperature colour bands",
], title_color=TEAL, bullet_size=10)


# ── Slide 28: GetCapabilities & DescribeSensor ────────────────────────────────
s = new_slide()
heading(s, "SOS GetCapabilities & DescribeSensor",
        "Clicking the Quick Request buttons proxies through FastAPI to the real istSOS server",
        bar_color=P2_BAR)

code_card(s, """\
<!-- SOS 2.0.0 GetCapabilities — from istSOS server -->

<sos:Capabilities version="2.0.0"
  xmlns:sos="http://www.opengis.net/sos/2.0">

  <ows:ServiceIdentification>
    <ows:Title>GeoPulse Sensor Observation Service</ows:Title>
    <ows:ServiceType>OGC:SOS</ows:ServiceType>
    <ows:ServiceTypeVersion>2.0.0</ows:ServiceTypeVersion>
  </ows:ServiceIdentification>

  <ows:OperationsMetadata>
    <ows:Operation name="GetObservation">
      <ows:DCP><ows:HTTP>
        <ows:Get xlink:href="/api/istsos/raw"/>
      </ows:HTTP></ows:DCP>
    </ows:Operation>
  </ows:OperationsMetadata>

  <sos:contents>
    <sos:ObservationOfferingList>
      <sos:ObservationOffering>
        <swes:identifier>urn:ogc:def:offering:geopulse</swes:identifier>
        <swes:procedure>Tokyo_Japan</swes:procedure>
        <swes:observableProperty>
          urn:ogc:def:parameter:x-istsos:1.0:meteo:air:temperature
        </swes:observableProperty>
      </sos:ObservationOffering>
    </sos:ObservationOfferingList>
  </sos:contents>
</sos:Capabilities>""",
    Inches(0.35), Inches(1.35), Inches(7.0), Inches(5.6))

card(s, Inches(7.35), Inches(1.35), Inches(5.63), Inches(2.7),
     "DescribeSensor Response", [
     "Procedure: Tokyo_Japan",
     "SensorML 2.0 output format",
     "Contains:",
     " Unique identifier string",
     " gml:Point — EPSG:4326 location",
     " OutputList — 5 observed properties",
     "   with OGC URN identifiers",
     " Units of measurement per output",
], title_color=TEAL)

card(s, Inches(7.35), Inches(4.25), Inches(5.63), Inches(2.7),
     "Frontend Quick Requests", [
     "Button → fetch(/api/istsos/capabilities)",
     "Raw XML shown in response panel",
     "DOMParser can parse it client-side",
     "",
     "DescribeSensor dropdown:",
     " select procedure name",
     " → fetch(/api/istsos/describe-sensor)",
     "   ?procedure=Tokyo_Japan",
], title_color=ACCENT)


# ── Slide 29: Deployment & Config ──────────────────────────────────────────────
s = new_slide()
heading(s, "Deployment — One-Command Startup",
        "docker compose up  handles database init, data loading, and API server automatically",
        bar_color=P2_BAR)

code_card(s, """\
# Full deployment sequence

# 1. Copy environment config
cp .env.example .env
# .env contains:
#   POSTGRES_DB=istsos  POSTGRES_USER=postgres
#   POSTGRES_PASSWORD=postgres  POSTGRES_HOST=postgis
#   APP_PORT=8000  APP_HOST=0.0.0.0
#   ISTSOS_SERVICE_URL=http://host.docker.internal:8088/istsos/geopulse
#   ISTSOS_DEFAULT_PROCEDURE=Tokyo_Japan

# 2. Start database + backend (background)
docker compose up -d postgis backend

# 3. Run one-shot data loader (waits for DB health)
docker compose run --rm loader
# → Loaded 17775 rows into geopulse.observations

# 4. Verify
curl http://127.0.0.1:8002/health
# → {"status":"ok","rows":17775,"istsos_service_url":"..."}

# 5. Open portal
#    http://127.0.0.1:8002

# --- To reload data ---
docker compose run --rm loader

# --- To shut down ---
docker compose down

# --- To reset database ---
docker compose down -v   # removes the postgis-data volume""",
    Inches(0.35), Inches(1.35), Inches(7.5), Inches(5.9))

card(s, Inches(8.0), Inches(1.35), Inches(4.98), Inches(2.8),
     "What Docker Gives Us", [
     "Reproducible environment — no 'works on",
     " my machine' — same image everywhere",
     "Service health checks — loader waits",
     " for postgis before inserting",
     "Volume persistence — data survives",
     " container restarts",
     "restart: no on loader — runs once,",
     " then exits cleanly",
], title_color=DOCKER)

card(s, Inches(8.0), Inches(4.35), Inches(4.98), Inches(2.9),
     "Backend Dockerfile", [
     "FROM python:3.11-slim",
     "COPY requirements.txt .",
     "RUN pip install -r requirements.txt",
     "  fastapi  uvicorn",
     "  psycopg2-binary  httpx",
     "COPY app/ /app/",
     "CMD uvicorn app.main:app",
     "  --host 0.0.0.0 --port 8000",
], title_color=TEAL)


# ── Slide 30: Summary & Standards Compliance ──────────────────────────────────
s = new_slide()
heading(s, "Project 2 — Summary & Standards Compliance",
        bar_color=P2_BAR)

stat_row(s, [
    ("25", "sensor stations"),
    ("17,775", "observations"),
    ("5", "OGC endpoints"),
    ("3", "Docker services"),
    ("1", "compose up"),
], Inches(1.35), P2_BAR)

data_table(s,
    ["OGC Standard", "Operation", "Implemented by", "Output"],
    [
        ["SOS 2.0", "GetCapabilities", "istSOS 3 (proxied)", "OGC XML"],
        ["SOS 2.0", "DescribeSensor",  "istSOS 3 (proxied)", "SensorML 2.0"],
        ["SOS 2.0", "GetObservation",  "FastAPI + PostGIS",  "O&M XML / JSON"],
        ["O&M",     "Spatial filter",  "ST_MakeEnvelope",   "GIST index scan"],
        ["SensorML","OGC URN IDs",     "prepare_import.py", "25 sensor defs"],
    ],
    Inches(0.35), Inches(2.65), Inches(12.63), Inches(2.8),
    col_widths=[Inches(2.0), Inches(2.6), Inches(4.2), Inches(3.83)]
)

card(s, Inches(0.35), Inches(5.65), Inches(6.0), Inches(1.55),
     "Upgrade from v1 (FastAPI + SQLite)", [
     "Real spatial DB (PostGIS GIST) vs flat SQLite",
     "Docker Compose — reproducible, portable",
     "OGC URN identifiers — standards-native",
     "istSOS proxy — path to full deployment",
], title_color=GREEN, bullet_size=10)

card(s, Inches(6.55), Inches(5.65), Inches(6.43), Inches(1.55),
     "Achieved", [
     "Fully containerised — one command startup",
     "Real OGC SOS 2.0 operations via istSOS",
     "Spatial BBOX queries via PostGIS geometry",
     "Interactive map + table + Chart.js charts",
], title_color=P2_BAR, bullet_size=10)


# ═════════════════════════════════════════════════════════════════════════════
prs.save("GeoPulse_istSOS_Presentation.pptx")
print(f"✅  Saved: GeoPulse_istSOS_Presentation.pptx  ({len(prs.slides)} slides)")
